import os
import logging
import hashlib
import subprocess
import tempfile
from pathlib import Path
import mimetypes
from typing import List, Dict, Any, Tuple, Optional
from datetime import datetime
import re
import shutil
import time

import pytesseract

logger = logging.getLogger(__name__)

pytesseract.pytesseract.tesseract_cmd = 'C:/Program Files/Tesseract-OCR/tesseract.exe'


class DocumentProcessor:
    """Universal document processor supporting multiple file types including Excel"""

    def __init__(self):
        self.supported_types = {
            'pdf': self._process_pdf,
            'docx': self._process_docx,
            'doc': self._process_doc,
            'txt': self._process_txt,
            'md': self._process_markdown,
            'rtf': self._process_rtf,
            'csv': self._process_csv,
            'xlsx': self._process_xlsx,
            'xlsm': self._process_xlsm,
            'xls': self._process_xls,
            'json': self._process_json,
            'xml': self._process_xml,
            'html': self._process_html,
            'htm': self._process_html,
            'pptx': self._process_pptx,
            'ppt': self._process_ppt,
            'jpg': self._process_image,
            'jpeg': self._process_image,
            'png': self._process_image,
            'bmp': self._process_image,
            'tiff': self._process_image,
            'tif': self._process_image,
            'gif': self._process_image,
            'webp': self._process_image
        }

        # Generic filenames that should be renamed based on content
        self.generic_filenames = {
            '1', '2', '3', '4', '5', '6', '7', '8', '9', '10',
            'document', 'file', 'untitled', 'new', 'temp', 'test',
            'doc', 'pdf', 'image', 'text', 'data', 'report',
            'copy', 'duplicate', 'unnamed', 'blank', 'sheet1',
            'workbook', 'book1', 'presentation', 'slide1'
        }

    def get_file_type(self, file_path: str) -> str:
        """Determine file type from file path"""
        return file_path.lower().split('.')[-1]

    def get_original_filename(self, file_path: str) -> str:
        """Extract original filename without extension"""
        return os.path.splitext(os.path.basename(file_path))[0]

    def is_generic_filename(self, filename: str) -> bool:
        """Check if filename is generic and should be renamed"""
        # Extract just the base name without extension
        base_name = os.path.splitext(os.path.basename(filename))[0].lower().strip()

        # Remove timestamp prefixes that might be added during upload
        original_base_name = re.sub(r'^\d{8}_\d{6}_', '', base_name)

        # If after removing timestamp, we have a meaningful name, it's not generic
        if len(original_base_name) > 2 and original_base_name not in self.generic_filenames:
            # Check if it's not just numbers
            if not original_base_name.isdigit():
                logger.debug(f"Filename '{original_base_name}' is considered valid (not generic)")
                return False

        # Check if it's in our generic list or is very short
        if original_base_name in self.generic_filenames or len(original_base_name) <= 2:
            logger.debug(f"Filename '{original_base_name}' is considered generic")
            return True

        # Check if it's just numbers
        if original_base_name.isdigit():
            logger.debug(f"Filename '{original_base_name}' is just numbers, considered generic")
            return True

        # Check for common generic patterns
        generic_patterns = [
            r'^untitled\d*$',
            r'^document\d*$',
            r'^file\d*$',
            r'^new\s*document\d*$',
            r'^scan\d*$',
            r'^img\d*$',
            r'^image\d*$',
            r'^sheet\d*$',
            r'^workbook\d*$',
            r'^book\d*$',
            r'^presentation\d*$'
        ]

        for pattern in generic_patterns:
            if re.match(pattern, original_base_name, re.IGNORECASE):
                logger.debug(f"Filename '{original_base_name}' matches generic pattern, considered generic")
                return True

        # If we get here, the filename appears to be meaningful
        logger.debug(f"Filename '{original_base_name}' is considered valid (not generic)")
        return False

    def is_supported(self, file_path: str) -> bool:
        """Check if file type is supported"""
        file_type = self.get_file_type(file_path)
        return file_type in self.supported_types

    def check_duplicate_by_hash_and_name(self, file_path: str, existing_files: List[Dict[str, Any]]) -> Optional[
        Dict[str, Any]]:
        """Enhanced duplicate checking"""
        current_hash = self.calculate_file_hash(file_path)
        current_name = self.get_clean_filename(file_path)

        for existing_file in existing_files:
            # Check both hash and cleaned filename
            if (existing_file.get('file_hash') == current_hash or
                    self.get_clean_filename(existing_file.get('filename', '')) == current_name):
                return existing_file

        return None

    def check_duplicate_by_hash(self, file_path: str, existing_files: List[Dict[str, Any]]) -> Optional[Dict[str, Any]]:
        """Check for duplicates by file hash"""
        current_hash = self.calculate_file_hash(file_path)

        for existing_file in existing_files:
            if existing_file.get('file_hash') == current_hash:
                return existing_file

        return None

    def process_document_atomic(self, file_path: str, **kwargs):
        """Atomic document processing to prevent duplicates"""
        temp_suffix = f".processing_{int(time.time())}"
        temp_path = f"{file_path}{temp_suffix}"

        try:
            # Move to temp location during processing
            shutil.move(file_path, temp_path)

            # Process the file
            result = self.process_document(temp_path, **kwargs)

            # Move to final location
            final_path = result[1]['file_path']
            if temp_path != final_path:
                shutil.move(temp_path, final_path)

            return result

        except Exception as e:
            # Restore original file on error
            if os.path.exists(temp_path):
                shutil.move(temp_path, file_path)
            raise e

    def can_convert_to_pdf(self, file_type: str) -> bool:
        """Check if file type can be converted to PDF"""
        convertible_types = ['docx', 'doc', 'xlsx', 'xlsm', 'xls', 'pptx', 'ppt', 'txt', 'rtf']
        return file_type.lower() in convertible_types

    def convert_to_pdf_for_viewing(self, file_path: str, output_dir: str = None) -> str:
        """Convert document to PDF for viewing purposes only"""
        try:
            file_type = self.get_file_type(file_path)

            if not self.can_convert_to_pdf(file_type):
                raise ValueError(f"Cannot convert {file_type} to PDF")

            # Create output directory if not specified
            if output_dir is None:
                output_dir = os.path.join(os.path.dirname(file_path), 'temp_pdfs')

            os.makedirs(output_dir, exist_ok=True)

            # Generate PDF filename
            base_name = os.path.splitext(os.path.basename(file_path))[0]
            pdf_filename = f"{base_name}_view.pdf"
            pdf_path = os.path.join(output_dir, pdf_filename)

            # Check if PDF already exists and is newer than source
            if os.path.exists(pdf_path):
                source_mtime = os.path.getmtime(file_path)
                pdf_mtime = os.path.getmtime(pdf_path)
                if pdf_mtime > source_mtime:
                    logger.info(f"Using existing PDF: {pdf_path}")
                    return pdf_path

            logger.info(f"Converting {file_type} to PDF: {file_path} -> {pdf_path}")

            if file_type.lower() in ['docx', 'doc']:
                success = self._convert_word_to_pdf(file_path, pdf_path)
            elif file_type.lower() in ['xlsx', 'xlsm', 'xls']:
                success = self._convert_excel_to_pdf(file_path, pdf_path)
            elif file_type.lower() in ['pptx', 'ppt']:
                success = self._convert_powerpoint_to_pdf(file_path, pdf_path)
            elif file_type.lower() == 'txt':
                success = self._convert_text_to_pdf(file_path, pdf_path)
            else:
                raise ValueError(f"No converter available for {file_type}")

            if success and os.path.exists(pdf_path):
                logger.info(f"Successfully converted to PDF: {pdf_path}")
                return pdf_path
            else:
                raise Exception("PDF conversion failed")

        except Exception as e:
            logger.error(f"Error converting {file_path} to PDF: {e}")
            raise

    def _convert_word_to_pdf(self, input_path: str, output_path: str) -> bool:
        """Convert Word document to PDF"""
        try:
            # Method 1: Try using docx2pdf (works on Windows with Word installed)
            try:
                from docx2pdf import convert
                convert(input_path, output_path)
                return os.path.exists(output_path)
            except ImportError:
                logger.warning("docx2pdf not available, trying alternative method")
            except Exception as e:
                logger.warning(f"docx2pdf failed: {e}, trying alternative method")

            # Method 2: Try using LibreOffice (cross-platform)
            try:
                result = subprocess.run([
                    'libreoffice', '--headless', '--convert-to', 'pdf',
                    '--outdir', os.path.dirname(output_path), input_path
                ], capture_output=True, text=True, timeout=60)

                if result.returncode == 0:
                    # LibreOffice creates PDF with same name as input
                    base_name = os.path.splitext(os.path.basename(input_path))[0]
                    libreoffice_pdf = os.path.join(os.path.dirname(output_path), f"{base_name}.pdf")

                    if os.path.exists(libreoffice_pdf):
                        if libreoffice_pdf != output_path:
                            os.rename(libreoffice_pdf, output_path)
                        return True

            except (subprocess.TimeoutExpired, FileNotFoundError) as e:
                logger.warning(f"LibreOffice conversion failed: {e}")

            # Method 3: Create PDF from extracted text (fallback)
            return self._create_pdf_from_text(input_path, output_path)

        except Exception as e:
            logger.error(f"Word to PDF conversion failed: {e}")
            return False

    def _convert_excel_to_pdf(self, input_path: str, output_path: str) -> bool:
        """Convert Excel document to PDF"""
        try:
            # Try LibreOffice first
            try:
                result = subprocess.run([
                    'libreoffice', '--headless', '--convert-to', 'pdf',
                    '--outdir', os.path.dirname(output_path), input_path
                ], capture_output=True, text=True, timeout=60)

                if result.returncode == 0:
                    base_name = os.path.splitext(os.path.basename(input_path))[0]
                    libreoffice_pdf = os.path.join(os.path.dirname(output_path), f"{base_name}.pdf")

                    if os.path.exists(libreoffice_pdf):
                        if libreoffice_pdf != output_path:
                            os.rename(libreoffice_pdf, output_path)
                        return True

            except (subprocess.TimeoutExpired, FileNotFoundError):
                logger.warning("LibreOffice not available for Excel conversion")

            # Fallback: Create PDF from extracted text
            return self._create_pdf_from_text(input_path, output_path)

        except Exception as e:
            logger.error(f"Excel to PDF conversion failed: {e}")
            return False

    def _convert_powerpoint_to_pdf(self, input_path: str, output_path: str) -> bool:
        """Convert PowerPoint document to PDF"""
        try:
            # Try LibreOffice first
            try:
                result = subprocess.run([
                    'libreoffice', '--headless', '--convert-to', 'pdf',
                    '--outdir', os.path.dirname(output_path), input_path
                ], capture_output=True, text=True, timeout=60)

                if result.returncode == 0:
                    base_name = os.path.splitext(os.path.basename(input_path))[0]
                    libreoffice_pdf = os.path.join(os.path.dirname(output_path), f"{base_name}.pdf")

                    if os.path.exists(libreoffice_pdf):
                        if libreoffice_pdf != output_path:
                            os.rename(libreoffice_pdf, output_path)
                        return True

            except (subprocess.TimeoutExpired, FileNotFoundError):
                logger.warning("LibreOffice not available for PowerPoint conversion")

            # Fallback: Create PDF from extracted text
            return self._create_pdf_from_text(input_path, output_path)

        except Exception as e:
            logger.error(f"PowerPoint to PDF conversion failed: {e}")
            return False

    def _convert_text_to_pdf(self, input_path: str, output_path: str) -> bool:
        """Convert text file to PDF"""
        return self._create_pdf_from_text(input_path, output_path)

    def _create_pdf_from_text(self, input_path: str, output_path: str) -> bool:
        """Create PDF from text content using reportlab"""
        try:
            from reportlab.lib.pagesizes import letter, A4
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.lib.units import inch

            # Extract text from the document
            file_type = self.get_file_type(input_path)
            processor = self.supported_types.get(file_type)

            if not processor:
                raise ValueError(f"No processor for {file_type}")

            text, metadata = processor(input_path)

            if not text:
                raise ValueError("No text content to convert")

            # Create PDF
            doc = SimpleDocTemplate(output_path, pagesize=A4)
            styles = getSampleStyleSheet()

            # Custom style for body text
            body_style = ParagraphStyle(
                'CustomBody',
                parent=styles['Normal'],
                fontSize=10,
                spaceAfter=12,
                leftIndent=0.5 * inch,
                rightIndent=0.5 * inch
            )

            story = []

            # Add title
            title = metadata.get('title') or os.path.splitext(os.path.basename(input_path))[0]
            story.append(Paragraph(title, styles['Title']))
            story.append(Spacer(1, 12))

            # Add metadata if available
            if metadata.get('author'):
                story.append(Paragraph(f"Author: {metadata['author']}", styles['Normal']))
            if metadata.get('subject'):
                story.append(Paragraph(f"Subject: {metadata['subject']}", styles['Normal']))

            story.append(Spacer(1, 20))

            # Add content
            paragraphs = text.split('\n\n')
            for para in paragraphs:
                if para.strip():
                    # Escape HTML characters and handle line breaks
                    para = para.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                    para = para.replace('\n', '<br/>')
                    story.append(Paragraph(para, body_style))
                    story.append(Spacer(1, 6))

            doc.build(story)
            return os.path.exists(output_path)

        except ImportError:
            logger.error("reportlab not installed. Install with: pip install reportlab")
            return False
        except Exception as e:
            logger.error(f"Text to PDF conversion failed: {e}")
            return False

    def process_document(self, file_path: str, chunk_size: int = 1000, chunk_overlap: int = 200,
                         auto_rename_generic: bool = True, max_title_length: int = 50,
                         existing_files: Optional[List[Dict[str, Any]]] = None,
                         original_filename: Optional[str] = None,
                         check_duplicates: bool = True) -> Tuple[List[Dict[str, Any]], Dict[str, Any]]:
        """Process any supported document type with enhanced filename handling"""

        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")

        file_type = self.get_file_type(file_path)

        if not self.is_supported(file_path):
            raise ValueError(f"Unsupported file type: {file_type}")

        try:
            logger.info(f"Processing {file_type.upper()} document: {file_path}")

            # Extract original filename info - use parameter if provided, otherwise extract from path
            original_filename = original_filename or os.path.basename(file_path)
            original_name_only = self.get_original_filename(original_filename)

            # Enhanced duplicate check BEFORE processing
            duplicate_info = None
            if check_duplicates and existing_files:
                duplicate_info = self.check_duplicate_by_hash_and_name(file_path, existing_files)
                if duplicate_info:
                    logger.warning(f"Duplicate file detected: {file_path}")
                    return [], {
                        "is_duplicate": True,
                        "duplicate_of": duplicate_info.get('id'),
                        "original_filename": original_filename,
                        "message": "File already exists in the system",
                        "existing_file": duplicate_info
                    }

            # Get processor function
            processor = self.supported_types[file_type]

            # Extract text and metadata
            text, metadata = processor(file_path)

            if not text or not text.strip():
                raise ValueError("No text content extracted from document")

            # Determine if we need to rename the file
            should_rename = False
            rename_reason = ""

            # Only rename if filename is actually generic AND auto_rename is enabled
            if auto_rename_generic:
                is_generic = self.is_generic_filename(original_filename)
                if is_generic:
                    should_rename = True
                    rename_reason = "generic_filename"
                    logger.info(f"Generic filename detected, will rename: {original_filename}")
                else:
                    logger.info(f"Valid filename detected, keeping original: {original_filename}")
            else:
                logger.info(f"Auto-rename disabled, keeping original filename: {original_filename}")

            # Rename file if needed with improved logic
            new_file_path = file_path
            if should_rename:
                new_file_path = self.rename_file_by_content(
                    file_path, text, metadata, max_title_length, rename_reason
                )
                logger.info(f"File renamed from: {original_filename} to: {os.path.basename(new_file_path)}")
            else:
                logger.info(f"File kept with original name: {original_filename}")

            # Create chunks
            chunks = self.create_chunks(text, chunk_size, chunk_overlap)

            # Enhanced processing metadata
            processing_metadata = {
                "original_filename": original_filename,
                "original_name_only": original_name_only,
                "final_filename": os.path.basename(new_file_path),
                "final_name_only": self.get_original_filename(new_file_path),
                "file_path": new_file_path,
                "file_type": file_type,
                "file_size": os.path.getsize(new_file_path),
                "file_hash": self.calculate_file_hash(new_file_path),
                "processing_date": datetime.utcnow().isoformat(),
                "chunk_count": len(chunks),
                "total_characters": len(text),
                "word_count": len(text.split()),
                "keywords": self.extract_keywords_simple(text),
                "was_renamed": new_file_path != file_path,
                "rename_reason": rename_reason if should_rename else None,
                "is_duplicate": False,
                "duplicate_of": None,
                "extracted_title": self.extract_title_from_content(text, metadata, file_type),
                "filename_was_generic": should_rename,
                "duplicate_check_enabled": check_duplicates,
                "processing_method": "enhanced_safe"
            }

            # Merge with file-specific metadata
            final_metadata = {**metadata, **processing_metadata}

            logger.info(f"Successfully processed {file_type.upper()}: {len(chunks)} chunks created")

            if processing_metadata["was_renamed"]:
                logger.info(f"File renamed: {original_filename} → {processing_metadata['final_filename']}")
            else:
                logger.info(f"File kept original name: {original_filename}")

            return chunks, final_metadata

        except Exception as e:
            logger.error(f"Error processing {file_type} document {file_path}: {e}")
            raise

    def rename_file_by_content(self, file_path: str, text: str, metadata: Dict[str, Any],
                               max_length: int = 50, reason: str = "generic_filename") -> str:
        """Rename file based on document title or content with improved duplicate handling"""
        try:
            file_type = self.get_file_type(file_path)
            file_dir = os.path.dirname(file_path)
            file_ext = os.path.splitext(file_path)[1]

            # Extract title based on file type
            title = self.extract_title_from_content(text, metadata, file_type)

            if not title:
                logger.info(f"No suitable title found for {file_path}, keeping original name")
                return file_path

            # Generate safe filename
            safe_filename = self.generate_safe_filename(title, max_length)
            new_filename = f"{safe_filename}{file_ext}"
            new_file_path = os.path.join(file_dir, new_filename)

            # If the new path is the same as original, no need to rename
            if os.path.normpath(new_file_path) == os.path.normpath(file_path):
                return file_path

            # Handle duplicates with timestamp instead of simple counter
            new_file_path = self.handle_duplicate_filename(new_file_path)

            # Perform the rename with retry logic for Windows
            max_retries = 3
            for attempt in range(max_retries):
                try:
                    if new_file_path != file_path:
                        # Use Windows-safe move operation
                        temp_suffix = f".tmp_{int(time.time())}"
                        temp_path = f"{file_path}{temp_suffix}"

                        # Move to temp first, then to final destination
                        os.rename(file_path, temp_path)
                        os.rename(temp_path, new_file_path)

                        logger.info(
                            f"Renamed file ({reason}): {os.path.basename(file_path)} → {os.path.basename(new_file_path)}")
                        return new_file_path
                    break

                except OSError as e:
                    if attempt < max_retries - 1:
                        time.sleep(0.1 * (attempt + 1))  # Progressive delay
                        # Try with different name
                        new_file_path = self.handle_duplicate_filename(new_file_path)
                    else:
                        logger.error(f"Failed to rename file after {max_retries} attempts: {e}")
                        return file_path

            return new_file_path

        except Exception as e:
            logger.error(f"Error renaming file {file_path}: {e}")
            return file_path

    def extract_title_from_content(self, text: str, metadata: Dict[str, Any], file_type: str) -> Optional[str]:
        """Extract title from document content based on file type"""

        if file_type == 'pdf':
            # Try metadata title first
            if metadata.get('title') and metadata['title'].strip():
                return metadata['title'].strip()
            # Fall back to first line of text
            return self.get_first_meaningful_line(text)

        elif file_type == 'docx':
            # Try document title property first
            if metadata.get('title') and metadata['title'].strip():
                return metadata['title'].strip()
            # Fall back to first paragraph
            return self.get_first_meaningful_line(text)

        elif file_type in ['doc', 'rtf']:
            # Use first meaningful line
            return self.get_first_meaningful_line(text)

        elif file_type == 'txt':
            # Use first line
            return self.get_first_meaningful_line(text)

        elif file_type == 'md':
            # Look for first header or title in front matter
            if 'title' in metadata and metadata['title']:
                return str(metadata['title']).strip()

            # Look for first markdown header
            header_match = re.search(r'^#+\s+(.+)', text, re.MULTILINE)
            if header_match:
                return header_match.group(1).strip()

            # Fall back to first line
            return self.get_first_meaningful_line(text)

        elif file_type == 'csv':
            # Use filename or create descriptive name from headers
            if metadata.get('headers'):
                headers = metadata['headers'][:3]  # First 3 headers
                return f"Data_{' '.join(headers)}"
            return None

        elif file_type in ['xlsx', 'xlsm', 'xls']:
            # Try workbook title or sheet names
            if metadata.get('title') and metadata['title'].strip():
                return metadata['title'].strip()

            # Use sheet names if available
            if metadata.get('sheet_names'):
                sheet_names = metadata['sheet_names']
                if len(sheet_names) == 1 and sheet_names[0] != 'Sheet1':
                    return f"Excel_{sheet_names[0]}"
                elif len(sheet_names) > 1:
                    return f"Excel_Workbook_{len(sheet_names)}_Sheets"

            # Fall back to first meaningful line from content
            return self.get_first_meaningful_line(text, max_words=6)

        elif file_type in ['pptx', 'ppt']:
            # Try presentation title
            if metadata.get('title') and metadata['title'].strip():
                return metadata['title'].strip()

            # Use first slide title if available
            if metadata.get('slide_titles') and metadata['slide_titles']:
                return metadata['slide_titles'][0]

            # Fall back to first meaningful line
            return self.get_first_meaningful_line(text)

        elif file_type == 'json':
            # Look for common title fields
            title_fields = ['title', 'name', 'filename', 'document_title']
            for field in title_fields:
                if field in metadata.get('structure', {}).get('keys', []):
                    # Extract from first few lines of text
                    lines = text.split('\n')[:10]
                    for line in lines:
                        if f"{field}:" in line:
                            title = line.split(':', 1)[1].strip()
                            if title:
                                return title
            return None

        elif file_type == 'xml':
            # Use root element name or first text content
            root_tag = metadata.get('root_tag', '').replace('{', '').replace('}', '')
            if root_tag:
                return f"{root_tag}_Document"
            return self.get_first_meaningful_line(text)

        elif file_type in ['html', 'htm']:
            # Use title tag first
            if metadata.get('title') and metadata['title'].strip():
                return metadata['title'].strip()

            # Look for first heading
            headings = metadata.get('headings', [])
            if headings:
                return headings[0].get('text', '').strip()

            # Fall back to first line
            return self.get_first_meaningful_line(text)

        elif file_type in ['jpg', 'jpeg', 'png', 'bmp', 'tiff', 'tif', 'gif', 'webp']:
            # Use first few words from OCR text
            return self.get_first_meaningful_line(text, max_words=8)

        return None

    def get_first_meaningful_line(self, text: str, max_words: int = 12) -> Optional[str]:
        """Extract first meaningful line from text"""
        if not text or not text.strip():
            return None

        lines = text.strip().split('\n')

        for line in lines:
            line = line.strip()

            # Skip empty lines, page markers, and very short lines
            if (len(line) < 3 or
                    line.startswith('[Page ') or
                    line.startswith('Page ') or
                    re.match(r'^\d+$', line) or  # Skip page numbers
                    len(line.split()) < 2):  # Skip single words
                continue

            # Clean up the line
            line = re.sub(r'\s+', ' ', line)  # Normalize whitespace
            line = re.sub(r'[^\w\s\-\.\,\!\?]', ' ', line)  # Remove special chars

            # Limit to max_words
            words = line.split()[:max_words]
            result = ' '.join(words)

            if len(result) >= 5:  # Minimum length
                return result

        return None

    def generate_safe_filename(self, title: str, max_length: int = 50) -> str:
        """Generate a safe filename from title"""
        if not title:
            return "untitled"

        # Remove or replace problematic characters
        # Keep only alphanumeric, spaces, hyphens, underscores, and periods
        safe_title = re.sub(r'[<>:"/\\|?*]', '', title)  # Remove illegal chars
        safe_title = re.sub(r'[^\w\s\-\.]', ' ', safe_title)  # Replace other special chars
        safe_title = re.sub(r'\s+', ' ', safe_title)  # Normalize spaces
        safe_title = safe_title.strip()

        # Replace spaces with underscores or hyphens
        safe_title = safe_title.replace(' ', '_')

        # Truncate to max length
        if len(safe_title) > max_length:
            safe_title = safe_title[:max_length].rstrip('_')

        # Ensure it's not empty
        if not safe_title:
            safe_title = "untitled"

        return safe_title

    def handle_duplicate_filename(self, file_path: str) -> str:
        """Handle duplicate filenames using timestamp instead of simple counter"""
        if not os.path.exists(file_path):
            return file_path

        base_path = os.path.splitext(file_path)[0]
        extension = os.path.splitext(file_path)[1]

        # Use timestamp for uniqueness instead of simple counter
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")

        # First try with timestamp
        new_path = f"{base_path}_{timestamp}{extension}"
        if not os.path.exists(new_path):
            return new_path

        # If timestamp exists, add milliseconds
        ms = int(time.time() * 1000) % 1000
        new_path = f"{base_path}_{timestamp}_{ms:03d}{extension}"

        # Final fallback with counter (should rarely be needed)
        counter = 1
        while os.path.exists(new_path) and counter <= 10:
            new_path = f"{base_path}_{timestamp}_{ms:03d}_{counter}{extension}"
            counter += 1

        return new_path

    def cleanup_numbered_duplicates(self, directory: str, dry_run: bool = True) -> Dict[str, Any]:
        """Clean up files with _1, _2, etc. suffixes that are duplicates"""
        results = {
            "scanned": 0,
            "duplicates_found": 0,
            "removed": 0,
            "kept": 0,
            "errors": 0,
            "files_processed": []
        }

        try:
            # Pattern to match files with numeric suffixes
            pattern = re.compile(r'^(.+)_(\d+)(\..+)$')

            for filename in os.listdir(directory):
                results["scanned"] += 1
                match = pattern.match(filename)

                if match:
                    base_name, number, extension = match.groups()
                    original_filename = f"{base_name}{extension}"
                    duplicate_path = os.path.join(directory, filename)
                    original_path = os.path.join(directory, original_filename)

                    if os.path.exists(original_path):
                        try:
                            # Compare file hashes
                            original_hash = self.calculate_file_hash(original_path)
                            duplicate_hash = self.calculate_file_hash(duplicate_path)

                            if original_hash == duplicate_hash:
                                results["duplicates_found"] += 1

                                if not dry_run:
                                    os.remove(duplicate_path)
                                    results["removed"] += 1
                                    logger.info(f"Removed duplicate: {filename}")
                                else:
                                    logger.info(f"Would remove duplicate: {filename}")

                                results["files_processed"].append({
                                    "file": filename,
                                    "action": "removed" if not dry_run else "would_remove",
                                    "reason": "identical_hash"
                                })
                            else:
                                results["kept"] += 1
                                results["files_processed"].append({
                                    "file": filename,
                                    "action": "kept",
                                    "reason": "different_content"
                                })

                        except Exception as e:
                            results["errors"] += 1
                            logger.error(f"Error processing {filename}: {e}")

            return results

        except Exception as e:
            logger.error(f"Error during cleanup: {e}")
            results["errors"] += 1
            return results

    def get_clean_filename(self, file_path: str) -> str:
        """Extract clean filename without date prefixes and extension"""
        filename = os.path.basename(file_path)
        name_only = os.path.splitext(filename)[0]

        # Remove date/timestamp prefixes that might have been added
        clean_name = re.sub(r'^\d{4}[-_]\d{2}[-_]\d{2}[-_]?', '', name_only)
        clean_name = re.sub(r'^\d{8}[-_]?', '', clean_name)
        clean_name = re.sub(r'^\d{2}[-_]\d{2}[-_]\d{4}[-_]?', '', clean_name)
        clean_name = re.sub(r'^\d{10,}[-_]?', '', clean_name)  # Remove timestamps
        clean_name = re.sub(r'^copy[-_]?of[-_]?', '', clean_name, flags=re.IGNORECASE)

        return clean_name.strip('_-') or name_only

    def cleanup_temp_files(self, temp_dir: str = None):
        """Clean up temporary files created during processing"""
        try:
            if temp_dir and os.path.exists(temp_dir):
                shutil.rmtree(temp_dir)
                logger.info(f"Cleaned up temporary directory: {temp_dir}")
        except Exception as e:
            logger.warning(f"Error cleaning up temp files: {e}")

    def clean_filename_for_storage(self, filename: str) -> str:
        """Clean filename for database storage (remove date prefixes)"""
        # Remove common date prefixes that might be added by the system
        clean_name = re.sub(r'^\d{4}[-_]\d{2}[-_]\d{2}[-_]?', '', filename)
        clean_name = re.sub(r'^\d{8}[-_]?', '', clean_name)
        clean_name = re.sub(r'^\d{2}[-_]\d{2}[-_]\d{4}[-_]?', '', clean_name)
        clean_name = re.sub(r'^\d{10,}[-_]?', '', clean_name)  # Remove timestamps

        return clean_name.strip('_-') or filename

    def get_display_filename(self, file_path: str, metadata: Dict[str, Any]) -> str:
        """Get clean filename for display purposes"""
        try:
            # Use the final filename from metadata if available
            if metadata.get('final_filename'):
                return self.clean_filename_for_storage(metadata['final_filename'])

            # Use original filename if available
            if metadata.get('original_filename'):
                return self.clean_filename_for_storage(metadata['original_filename'])

            # Fall back to extracting from file path
            return self.clean_filename_for_storage(os.path.basename(file_path))

        except Exception as e:
            logger.error(f"Error getting display filename: {e}")
            return os.path.basename(file_path)

    def get_display_name(self, metadata: Dict[str, Any]) -> str:
        """Get the best display name for the document"""
        try:
            # For files that weren't renamed (valid filenames), use original filename
            if not metadata.get('was_renamed', False) and metadata.get('original_filename'):
                display_name = os.path.splitext(metadata['original_filename'])[0]

            # For renamed files, use extracted title or final filename
            elif metadata.get('extracted_title'):
                display_name = metadata['extracted_title']

            elif metadata.get('final_name_only'):
                display_name = metadata['final_name_only']

            elif metadata.get('original_name_only'):
                display_name = metadata['original_name_only']

            else:
                display_name = 'Untitled Document'

            # Clean up display name
            display_name = display_name.strip()

            # Ensure it's not too long for display
            if len(display_name) > 60:
                display_name = display_name[:57] + "..."

            return display_name

        except Exception as e:
            logger.error(f"Error getting display name: {e}")
            return 'Untitled Document'

    def _process_pdf(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        """Process PDF files"""
        try:
            import PyPDF2

            text = ""
            metadata = {}

            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)

                # Extract metadata
                if pdf_reader.metadata:
                    metadata.update({
                        "title": pdf_reader.metadata.get('/Title', ''),
                        "author": pdf_reader.metadata.get('/Author', ''),
                        "subject": pdf_reader.metadata.get('/Subject', ''),
                        "creator": pdf_reader.metadata.get('/Creator', ''),
                        "producer": pdf_reader.metadata.get('/Producer', ''),
                        "creation_date": str(pdf_reader.metadata.get('/CreationDate', '')),
                        "modification_date": str(pdf_reader.metadata.get('/ModDate', ''))
                    })

                metadata["total_pages"] = len(pdf_reader.pages)

                # Extract text from all pages
                for page_num, page in enumerate(pdf_reader.pages):
                    try:
                        page_text = page.extract_text()
                        if page_text:
                            text += f"\n[Page {page_num + 1}]\n{page_text}\n"
                    except Exception as e:
                        logger.warning(f"Could not extract text from page {page_num + 1}: {e}")

            return text.strip(), metadata

        except ImportError:
            raise ImportError("PyPDF2 is required for PDF processing. Install with: pip install PyPDF2")
        except Exception as e:
            logger.error(f"Error processing PDF: {e}")
            raise

    def _process_docx(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        """Process DOCX files"""
        try:
            from docx import Document

            doc = Document(file_path)

            # Extract text
            text = ""
            paragraph_count = 0

            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text += paragraph.text + "\n"
                    paragraph_count += 1

            # Extract metadata
            metadata = {
                "title": doc.core_properties.title or "",
                "author": doc.core_properties.author or "",
                "subject": doc.core_properties.subject or "",
                "created": str(doc.core_properties.created) if doc.core_properties.created else "",
                "modified": str(doc.core_properties.modified) if doc.core_properties.modified else "",
                "paragraph_count": paragraph_count
            }

            return text.strip(), metadata

        except ImportError:
            raise ImportError(
                "python-docx is required for DOCX processing. Install with: pip install python-docx")
        except Exception as e:
            logger.error(f"Error processing DOCX: {e}")
            raise

    def _process_doc(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        """Process DOC files (legacy Word format)"""

        # Method 1: Try using textract with antiword
        try:
            import textract
            text = textract.process(file_path).decode('utf-8')

            metadata = {
                "format": "doc",
                "extracted_with": "textract"
            }
            return text.strip(), metadata

        except ImportError:
            logger.warning("textract not available for DOC processing")
        except Exception as e:
            logger.warning(f"textract failed for DOC processing: {e}")

        # Method 2: Try using LibreOffice headless conversion
        try:
            # Convert DOC to TXT using LibreOffice
            temp_dir = tempfile.mkdtemp()
            try:
                result = subprocess.run([
                    'libreoffice', '--headless', '--convert-to', 'txt',
                    '--outdir', temp_dir, file_path
                ], capture_output=True, text=True, timeout=60)

                if result.returncode == 0:
                    # Find the converted txt file
                    base_name = os.path.splitext(os.path.basename(file_path))[0]
                    txt_file = os.path.join(temp_dir, f"{base_name}.txt")

                    if os.path.exists(txt_file):
                        with open(txt_file, 'r', encoding='utf-8') as f:
                            text = f.read()

                        metadata = {
                            "format": "doc",
                            "extracted_with": "libreoffice"
                        }
                        return text.strip(), metadata
            finally:
                shutil.rmtree(temp_dir, ignore_errors=True)

        except (subprocess.TimeoutExpired, FileNotFoundError) as e:
            logger.warning(f"LibreOffice conversion failed: {e}")
        except Exception as e:
            logger.warning(f"LibreOffice method failed: {e}")

        # Method 3: Try using python-docx (limited compatibility with .doc)
        try:
            from docx import Document

            # This might work for some .doc files that are actually .docx
            doc = Document(file_path)
            text = ""

            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text += paragraph.text + "\n"

            if text.strip():
                metadata = {
                    "format": "doc",
                    "extracted_with": "python-docx",
                    "note": "Limited compatibility - some formatting may be lost"
                }
                return text.strip(), metadata

        except Exception as e:
            logger.warning(f"python-docx method failed: {e}")

            # Method 4: Try using win32com (Windows only)
        if os.name == 'nt':  # Windows
            try:
                import win32com.client

                word = win32com.client.Dispatch("Word.Application")
                word.Visible = False

                try:
                    doc = word.Documents.Open(file_path)
                    text = doc.Content.Text
                    doc.Close()

                    metadata = {
                        "format": "doc",
                        "extracted_with": "win32com",
                        "note": "Extracted using Microsoft Word COM interface"
                    }
                    return text.strip(), metadata

                finally:
                    word.Quit()

            except ImportError:
                logger.warning("pywin32 not available for DOC processing")
            except Exception as e:
                logger.warning(f"Win32COM method failed: {e}")

            # Method 5: Last resort - try to read as RTF (some .doc files are RTF)
        try:
            from striprtf.striprtf import rtf_to_text

            with open(file_path, 'r', encoding='latin-1', errors='ignore') as file:
                content = file.read()

            # Check if it might be RTF format
            if content.startswith('{\\rtf'):
                text = rtf_to_text(content)

                metadata = {
                    "format": "doc",
                    "extracted_with": "striprtf",
                    "note": "Processed as RTF format"
                }
                return text.strip(), metadata

        except Exception as e:
            logger.warning(f"RTF fallback method failed: {e}")

            # If all methods fail, raise an informative error
        raise ValueError(
            f"Unable to process DOC file: {file_path}. "
            "Please install one of the following:\n"
            "1. LibreOffice (recommended): https://www.libreoffice.org/\n"
            "2. pywin32 (Windows): pip install pywin32\n"
            "3. textract with antiword: pip install textract (requires antiword installation)"
        )

    def _process_xlsx(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        """Process XLSX files (Excel 2007+)"""
        try:
            import openpyxl
            from openpyxl.utils import get_column_letter

            workbook = openpyxl.load_workbook(file_path, data_only=True)

            text_parts = []
            metadata = {
                "format": "xlsx",
                "sheet_names": workbook.sheetnames,
                "sheet_count": len(workbook.sheetnames),
                "total_rows": 0,
                "total_cells_with_data": 0
            }

            # Extract workbook properties if available
            if hasattr(workbook, 'properties') and workbook.properties:
                props = workbook.properties
                metadata.update({
                    "title": getattr(props, 'title', '') or '',
                    "author": getattr(props, 'creator', '') or '',
                    "subject": getattr(props, 'subject', '') or '',
                    "created": str(getattr(props, 'created', '')) if getattr(props, 'created',
                                                                             None) else '',
                    "modified": str(getattr(props, 'modified', '')) if getattr(props, 'modified',
                                                                               None) else ''
                })

            # Process each worksheet
            for sheet_name in workbook.sheetnames:
                try:
                    worksheet = workbook[sheet_name]

                    # Add sheet header
                    text_parts.append(f"\n=== Sheet: {sheet_name} ===\n")

                    # Get the used range
                    if worksheet.max_row == 1 and worksheet.max_column == 1:
                        # Empty sheet
                        text_parts.append("(Empty sheet)")
                        continue

                    # Extract data row by row
                    sheet_rows = 0
                    sheet_cells_with_data = 0

                    for row_num, row in enumerate(worksheet.iter_rows(values_only=True), 1):
                        # Skip completely empty rows
                        if not any(cell is not None and str(cell).strip() for cell in row):
                            continue

                        sheet_rows += 1
                        row_text = []

                        for col_num, cell_value in enumerate(row, 1):
                            if cell_value is not None and str(cell_value).strip():
                                col_letter = get_column_letter(col_num)

                                # Format cell value
                                if isinstance(cell_value, (int, float)):
                                    formatted_value = str(cell_value)
                                else:
                                    formatted_value = str(cell_value).strip()

                                if formatted_value:
                                    row_text.append(f"{col_letter}{row_num}: {formatted_value}")
                                    sheet_cells_with_data += 1

                        if row_text:
                            text_parts.append(" | ".join(row_text))

                        # Limit rows per sheet to prevent huge documents
                        if sheet_rows > 500:
                            text_parts.append(f"... (truncated at 500 rows for sheet '{sheet_name}')")
                            break

                    metadata["total_rows"] += sheet_rows
                    metadata["total_cells_with_data"] += sheet_cells_with_data

                except Exception as e:
                    logger.warning(f"Error processing sheet '{sheet_name}': {e}")
                    text_parts.append(f"Error processing sheet '{sheet_name}': {str(e)}")

            text = "\n".join(text_parts)

            # Add summary to metadata
            metadata[
                "processing_summary"] = f"Processed {len(workbook.sheetnames)} sheets with {metadata['total_rows']} total rows"

            return text, metadata

        except ImportError:
            raise ImportError(
                "openpyxl is required for XLSX processing. Install with: pip install openpyxl")
        except Exception as e:
            logger.error(f"Error processing XLSX: {e}")
            raise

    def _process_xlsm(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        """Process XLSM files (Excel 2007+ with macros)"""
        try:
            import openpyxl
            from openpyxl.utils import get_column_letter

            # Load workbook with data_only=True to get calculated values, keep_vba=False since we only need data
            workbook = openpyxl.load_workbook(file_path, data_only=True, keep_vba=False)

            text_parts = []
            metadata = {
                "format": "xlsm",
                "contains_macros": True,  # XLSM files always contain macros
                "sheet_names": workbook.sheetnames,
                "sheet_count": len(workbook.sheetnames),
                "total_rows": 0,
                "total_cells_with_data": 0
            }

            # Extract workbook properties if available
            if hasattr(workbook, 'properties') and workbook.properties:
                props = workbook.properties
                metadata.update({
                    "title": getattr(props, 'title', '') or '',
                    "author": getattr(props, 'creator', '') or '',
                    "subject": getattr(props, 'subject', '') or '',
                    "created": str(getattr(props, 'created', '')) if getattr(props, 'created',
                                                                             None) else '',
                    "modified": str(getattr(props, 'modified', '')) if getattr(props, 'modified',
                                                                               None) else ''
                })

            # Add macro warning to text
            text_parts.append("=== EXCEL FILE WITH MACROS (XLSM) ===")
            text_parts.append(
                "Note: This file contains macros. Only data values are extracted, not macro code.\n")

            # Process each worksheet
            for sheet_name in workbook.sheetnames:
                try:
                    worksheet = workbook[sheet_name]

                    # Add sheet header
                    text_parts.append(f"\n=== Sheet: {sheet_name} ===\n")

                    # Get the used range
                    if worksheet.max_row == 1 and worksheet.max_column == 1:
                        # Empty sheet
                        text_parts.append("(Empty sheet)")
                        continue

                    # Extract data row by row
                    sheet_rows = 0
                    sheet_cells_with_data = 0

                    for row_num, row in enumerate(worksheet.iter_rows(values_only=True), 1):
                        # Skip completely empty rows
                        if not any(cell is not None and str(cell).strip() for cell in row):
                            continue

                        sheet_rows += 1
                        row_text = []

                        for col_num, cell_value in enumerate(row, 1):
                            if cell_value is not None and str(cell_value).strip():
                                col_letter = get_column_letter(col_num)

                                # Format cell value (calculated values for formula cells)
                                if isinstance(cell_value, (int, float)):
                                    formatted_value = str(cell_value)
                                else:
                                    formatted_value = str(cell_value).strip()

                                if formatted_value:
                                    row_text.append(f"{col_letter}{row_num}: {formatted_value}")
                                    sheet_cells_with_data += 1

                        if row_text:
                            text_parts.append(" | ".join(row_text))

                        # Limit rows per sheet to prevent huge documents
                        if sheet_rows > 500:
                            text_parts.append(f"... (truncated at 500 rows for sheet '{sheet_name}')")
                            break

                    metadata["total_rows"] += sheet_rows
                    metadata["total_cells_with_data"] += sheet_cells_with_data

                except Exception as e:
                    logger.warning(f"Error processing sheet '{sheet_name}': {e}")
                    text_parts.append(f"Error processing sheet '{sheet_name}': {str(e)}")

            text = "\n".join(text_parts)

            # Add summary to metadata
            metadata[
                "processing_summary"] = f"Processed {len(workbook.sheetnames)} sheets with {metadata['total_rows']} total rows (macro-enabled file)"
            metadata[
                "macro_warning"] = "This file contains macros. Only data values were extracted for security."

            return text, metadata

        except ImportError:
            raise ImportError(
                "openpyxl is required for XLSM processing. Install with: pip install openpyxl")
        except Exception as e:
            logger.error(f"Error processing XLSM: {e}")
            raise

    def _process_xls(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        """Process XLS files (legacy Excel format)"""
        try:
            import xlrd

            workbook = xlrd.open_workbook(file_path)

            text_parts = []
            metadata = {
                "format": "xls",
                "sheet_names": workbook.sheet_names(),
                "sheet_count": workbook.nsheets,
                "total_rows": 0,
                "total_cells_with_data": 0
            }

            # Process each worksheet
            for sheet_name in workbook.sheet_names():
                try:
                    worksheet = workbook.sheet_by_name(sheet_name)

                    # Add sheet header
                    text_parts.append(f"\n=== Sheet: {sheet_name} ===\n")

                    if worksheet.nrows == 0:
                        text_parts.append("(Empty sheet)")
                        continue

                    # Extract data row by row
                    sheet_rows = 0
                    sheet_cells_with_data = 0

                    for row_num in range(worksheet.nrows):
                        row_values = worksheet.row_values(row_num)

                        # Skip completely empty rows
                        if not any(str(cell).strip() for cell in row_values if cell):
                            continue

                        sheet_rows += 1
                        row_text = []

                        for col_num, cell_value in enumerate(row_values):
                            if cell_value and str(cell_value).strip():
                                # Convert column number to letter
                                col_letter = chr(
                                    65 + col_num) if col_num < 26 else f"A{chr(65 + col_num - 26)}"

                                # Format cell value
                                if isinstance(cell_value, (int, float)):
                                    formatted_value = str(cell_value)
                                else:
                                    formatted_value = str(cell_value).strip()

                                if formatted_value:
                                    row_text.append(f"{col_letter}{row_num + 1}: {formatted_value}")
                                    sheet_cells_with_data += 1

                        if row_text:
                            text_parts.append(" | ".join(row_text))

                        # Limit rows per sheet to prevent huge documents
                        if sheet_rows > 500:
                            text_parts.append(f"... (truncated at 500 rows for sheet '{sheet_name}')")
                            break

                    metadata["total_rows"] += sheet_rows
                    metadata["total_cells_with_data"] += sheet_cells_with_data

                except Exception as e:
                    logger.warning(f"Error processing sheet '{sheet_name}': {e}")
                    text_parts.append(f"Error processing sheet '{sheet_name}': {str(e)}")

            text = "\n".join(text_parts)

            # Add summary to metadata
            metadata[
                "processing_summary"] = f"Processed {workbook.nsheets} sheets with {metadata['total_rows']} total rows"

            return text, metadata

        except ImportError:
            raise ImportError("xlrd is required for XLS processing. Install with: pip install xlrd")
        except Exception as e:
            logger.error(f"Error processing XLS: {e}")
            raise

    def _process_pptx(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        """Process PPTX files (PowerPoint 2007+)"""
        try:
            from pptx import Presentation

            presentation = Presentation(file_path)

            text_parts = []
            slide_titles = []
            metadata = {
                "format": "pptx",
                "slide_count": len(presentation.slides),
                "total_shapes": 0,
                "total_text_shapes": 0
            }

            # Extract presentation properties if available
            if hasattr(presentation, 'core_properties') and presentation.core_properties:
                props = presentation.core_properties
                metadata.update({
                    "title": getattr(props, 'title', '') or '',
                    "author": getattr(props, 'author', '') or '',
                    "subject": getattr(props, 'subject', '') or '',
                    "created": str(getattr(props, 'created', '')) if getattr(props, 'created', None) else '',
                    "modified": str(getattr(props, 'modified', '')) if getattr(props, 'modified', None) else ''
                })

            # Process each slide
            for slide_num, slide in enumerate(presentation.slides, 1):
                try:
                    text_parts.append(f"\n=== Slide {slide_num} ===\n")

                    slide_text = []
                    slide_title = None

                    # Extract text from all shapes
                    for shape in slide.shapes:
                        metadata["total_shapes"] += 1

                        if hasattr(shape, "text") and shape.text.strip():
                            metadata["total_text_shapes"] += 1
                            shape_text = shape.text.strip()

                            # First text shape is often the title
                            if slide_title is None and len(shape_text) < 100:
                                slide_title = shape_text
                                slide_titles.append(shape_text)

                            slide_text.append(shape_text)

                    if slide_text:
                        text_parts.extend(slide_text)
                    else:
                        text_parts.append("(No text content)")

                except Exception as e:
                    logger.warning(f"Error processing slide {slide_num}: {e}")
                    text_parts.append(f"Error processing slide {slide_num}: {str(e)}")

            text = "\n".join(text_parts)
            metadata["slide_titles"] = slide_titles
            metadata[
                "processing_summary"] = f"Processed {len(presentation.slides)} slides with {metadata['total_text_shapes']} text shapes"

            return text, metadata

        except ImportError:
            raise ImportError("python-pptx is required for PPTX processing. Install with: pip install python-pptx")
        except Exception as e:
            logger.error(f"Error processing PPTX: {e}")
            raise

    def _process_ppt(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        """Process PPT files (legacy PowerPoint format)"""
        try:
            import textract

            text = textract.process(file_path).decode('utf-8')

            metadata = {
                "format": "ppt",
                "extracted_with": "textract",
                "processing_summary": "Processed legacy PowerPoint file using textract"
            }

            return text.strip(), metadata

        except ImportError:
            raise ImportError("textract is required for PPT processing. Install with: pip install textract")
        except Exception as e:
            logger.error(f"Error processing PPT: {e}")
            raise

    def _process_txt(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        """Process plain text files"""
        try:
            # Try different encodings
            encodings = ['utf-8', 'utf-16', 'latin-1', 'cp1252']
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as file:
                        text = file.read()
                    break
                except UnicodeDecodeError:
                    continue
            else:
                raise ValueError("Could not decode text file with any supported encoding")

            metadata = {
                "format": "txt",
                "encoding": encoding,
                "line_count": len(text.splitlines()),
                "processing_summary": f"Processed text file with {encoding} encoding"
            }

            return text, metadata

        except Exception as e:
            logger.error(f"Error processing TXT: {e}")
            raise

    def _process_markdown(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        """Process Markdown files"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                text = file.read()

            # Extract markdown metadata (front matter)
            metadata = {"format": "markdown"}

            # Check for YAML front matter
            if text.startswith('---'):
                try:
                    import yaml
                    parts = text.split('---', 2)
                    if len(parts) >= 3:
                        front_matter = yaml.safe_load(parts[1])
                        if isinstance(front_matter, dict):
                            metadata.update(front_matter)
                        text = parts[2].strip()
                except ImportError:
                    logger.warning("PyYAML not available for front matter parsing")
                except Exception as e:
                    logger.warning(f"Could not parse front matter: {e}")

            # Count headers
            headers = re.findall(r'^#+\s+(.+)', text, re.MULTILINE)
            metadata["header_count"] = len(headers)
            metadata["headers"] = headers[:10]  # First 10 headers
            metadata["processing_summary"] = f"Processed Markdown file with {len(headers)} headers"

            return text, metadata

        except Exception as e:
            logger.error(f"Error processing Markdown: {e}")
            raise

    def _process_rtf(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        """Process RTF files"""
        try:
            from striprtf.striprtf import rtf_to_text

            with open(file_path, 'r', encoding='utf-8') as file:
                rtf_content = file.read()

            text = rtf_to_text(rtf_content)

            metadata = {
                "format": "rtf",
                "extracted_with": "striprtf",
                "processing_summary": "Processed RTF file using striprtf"
            }

            return text, metadata

        except ImportError:
            raise ImportError("striprtf is required for RTF processing. Install with: pip install striprtf")
        except Exception as e:
            logger.error(f"Error processing RTF: {e}")
            raise

    def _process_csv(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        """Process CSV files"""
        try:
            import csv

            text_parts = []
            row_count = 0
            headers = []

            with open(file_path, 'r', encoding='utf-8') as file:
                # Detect delimiter
                sample = file.read(1024)
                file.seek(0)
                sniffer = csv.Sniffer()
                delimiter = sniffer.sniff(sample).delimiter

                reader = csv.reader(file, delimiter=delimiter)

                for i, row in enumerate(reader):
                    if i == 0:
                        headers = row
                        text_parts.append("Headers: " + ", ".join(row))
                    else:
                        # Convert row to readable text
                        row_text = []
                        for j, cell in enumerate(row):
                            if j < len(headers):
                                row_text.append(f"{headers[j]}: {cell}")
                            else:
                                row_text.append(cell)
                        text_parts.append("Row " + str(i) + ": " + "; ".join(row_text))

                    row_count += 1

                    # Limit rows to prevent huge documents
                    if row_count > 1000:
                        text_parts.append(f"... (truncated at 1000 rows, total rows: {row_count})")
                        break

            text = "\n".join(text_parts)

            metadata = {
                "format": "csv",
                "delimiter": delimiter,
                "headers": headers,
                "row_count": row_count,
                "column_count": len(headers),
                "processing_summary": f"Processed CSV file with {row_count} rows and {len(headers)} columns"
            }

            return text, metadata

        except Exception as e:
            logger.error(f"Error processing CSV: {e}")
            raise

    def _process_json(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        """Process JSON files"""
        try:
            import json

            with open(file_path, 'r', encoding='utf-8') as file:
                data = json.load(file)

            # Convert JSON to readable text
            def json_to_text(obj, prefix=""):
                text_parts = []
                if isinstance(obj, dict):
                    for key, value in obj.items():
                        current_prefix = f"{prefix}.{key}" if prefix else key
                        if isinstance(value, (dict, list)):
                            text_parts.extend(json_to_text(value, current_prefix))
                        else:
                            text_parts.append(f"{current_prefix}: {value}")
                elif isinstance(obj, list):
                    for i, item in enumerate(obj):
                        current_prefix = f"{prefix}[{i}]" if prefix else f"item_{i}"
                        if isinstance(item, (dict, list)):
                            text_parts.extend(json_to_text(item, current_prefix))
                        else:
                            text_parts.append(f"{current_prefix}: {item}")
                return text_parts

            text_parts = json_to_text(data)
            text = "\n".join(text_parts)

            # Analyze JSON structure
            def analyze_structure(obj):
                if isinstance(obj, dict):
                    return {
                        "type": "object",
                        "keys": list(obj.keys())[:20],  # First 20 keys
                        "key_count": len(obj.keys())
                    }
                elif isinstance(obj, list):
                    return {
                        "type": "array",
                        "length": len(obj),
                        "item_types": list(set(type(item).__name__ for item in obj[:10]))
                    }
                else:
                    return {"type": type(obj).__name__}

            metadata = {
                "format": "json",
                "structure": analyze_structure(data),
                "file_size_bytes": os.path.getsize(file_path),
                "processing_summary": f"Processed JSON file with {len(text_parts)} data elements"
            }

            return text, metadata

        except Exception as e:
            logger.error(f"Error processing JSON: {e}")
            raise

    def _process_xml(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        """Process XML files"""
        try:
            import xml.etree.ElementTree as ET

            tree = ET.parse(file_path)
            root = tree.getroot()

            # Convert XML to readable text
            def xml_to_text(element, level=0):
                text_parts = []
                indent = "  " * level

                # Element name and attributes
                elem_text = f"{indent}{element.tag}"
                if element.attrib:
                    attrs = ", ".join(f"{k}={v}" for k, v in element.attrib.items())
                    elem_text += f" ({attrs})"

                # Element text content
                if element.text and element.text.strip():
                    elem_text += f": {element.text.strip()}"

                text_parts.append(elem_text)

                # Process children
                for child in element:
                    text_parts.extend(xml_to_text(child, level + 1))

                return text_parts

            text_parts = xml_to_text(root)
            text = "\n".join(text_parts)

            # Count elements
            def count_elements(element):
                count = 1
                for child in element:
                    count += count_elements(child)
                return count

            metadata = {
                "format": "xml",
                "root_tag": root.tag,
                "total_elements": count_elements(root),
                "root_attributes": dict(root.attrib) if root.attrib else {},
                "processing_summary": f"Processed XML file with {count_elements(root)} elements"
            }

            return text, metadata

        except Exception as e:
            logger.error(f"Error processing XML: {e}")
            raise

    def _process_html(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        """Process HTML files"""
        try:
            from bs4 import BeautifulSoup

            with open(file_path, 'r', encoding='utf-8') as file:
                html_content = file.read()

            soup = BeautifulSoup(html_content, 'html.parser')

            # Extract text content
            text = soup.get_text(separator='\n', strip=True)

            # Extract metadata
            metadata = {
                "format": "html",
                "title": soup.title.string if soup.title else "",
                "meta_description": "",
                "meta_keywords": "",
                "links_count": len(soup.find_all('a')),
                "images_count": len(soup.find_all('img')),
                "headings": []
            }

            # Extract meta tags
            meta_desc = soup.find('meta', attrs={'name': 'description'})
            if meta_desc:
                metadata["meta_description"] = meta_desc.get('content', '')

            meta_keywords = soup.find('meta', attrs={'name': 'keywords'})
            if meta_keywords:
                metadata["meta_keywords"] = meta_keywords.get('content', '')

            # Extract headings
            for i in range(1, 7):
                headings = soup.find_all(f'h{i}')
                for heading in headings[:5]:  # First 5 headings of each level
                    metadata["headings"].append({
                        "level": i,
                        "text": heading.get_text(strip=True)
                    })

            metadata["processing_summary"] = f"Processed HTML file with {len(metadata['headings'])} headings and {metadata['links_count']} links"

            return text, metadata

        except ImportError:
            raise ImportError("beautifulsoup4 is required for HTML processing. Install with: pip install beautifulsoup4")
        except Exception as e:
            logger.error(f"Error processing HTML: {e}")
            raise

    def _process_image(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        """Process image files using OCR"""
        try:
            from PIL import Image
            import pytesseract

            # Open image
            image = Image.open(file_path)

            # Perform OCR
            ocr_config = '--oem 3 --psm 6'  # OCR Engine Mode 3, Page Segmentation Mode 6

            # Get text with confidence scores
            ocr_data = pytesseract.image_to_data(image, config=ocr_config, output_type=pytesseract.Output.DICT)

            # Filter text by confidence
            min_confidence = 30  # Minimum confidence threshold
            text_parts = []
            confidences = []

            for i, conf in enumerate(ocr_data['conf']):
                if int(conf) > min_confidence:
                    text = ocr_data['text'][i].strip()
                    if text:
                        text_parts.append(text)
                        confidences.append(int(conf))

            extracted_text = ' '.join(text_parts)

            # Image metadata
            metadata = {
                "format": "image",
                "image_format": image.format,
                "size": image.size,
                "mode": image.mode,
                                "ocr_confidence_avg": sum(confidences) / len(confidences) if confidences else 0,
                "ocr_confidence_min": min(confidences) if confidences else 0,
                "ocr_confidence_max": max(confidences) if confidences else 0,
                "ocr_words_detected": len(text_parts),
                "has_alpha": image.mode in ('RGBA', 'LA'),
                "dpi": image.info.get('dpi', (72, 72)),
                "processing_summary": f"Processed {image.format} image with OCR, extracted {len(text_parts)} words"
            }

            # Add EXIF data if available
            if hasattr(image, '_getexif') and image._getexif():
                try:
                    exif = image._getexif()
                    metadata["exif"] = {k: str(v) for k, v in exif.items() if isinstance(v, (str, int, float))}
                except:
                    pass

            return extracted_text, metadata

        except ImportError:
            raise ImportError("PIL and pytesseract are required for image processing. Install with: pip install Pillow pytesseract")
        except Exception as e:
            logger.error(f"Error processing image: {e}")
            raise

    def create_chunks(self, text: str, chunk_size: int = 1000, chunk_overlap: int = 200) -> List[Dict[str, Any]]:
        """Create text chunks with metadata"""
        if not text or not text.strip():
            return []

        # Split text into sentences for better chunking
        sentences = self.simple_sentence_split(text)
        chunks = []
        current_chunk = ""
        current_sentences = []
        chunk_index = 0

        for sentence in sentences:
            # Check if adding this sentence would exceed chunk size
            if len(current_chunk) + len(sentence) > chunk_size and current_chunk:
                # Create chunk
                chunk_data = {
                    "chunk_index": chunk_index,
                    "content": current_chunk.strip(),
                    "word_count": len(current_chunk.split()),
                    "char_count": len(current_chunk),
                    "sentence_count": len(current_sentences)
                }
                chunks.append(chunk_data)

                # Handle overlap
                if chunk_overlap > 0 and current_sentences:
                    # Keep last few sentences for overlap
                    overlap_text = ""
                    overlap_sentences = []
                    for sent in reversed(current_sentences):
                        if len(overlap_text) + len(sent) <= chunk_overlap:
                            overlap_text = sent + " " + overlap_text
                            overlap_sentences.insert(0, sent)
                        else:
                            break

                    current_chunk = overlap_text + sentence
                    current_sentences = overlap_sentences + [sentence]
                else:
                    current_chunk = sentence
                    current_sentences = [sentence]

                chunk_index += 1
            else:
                current_chunk += " " + sentence if current_chunk else sentence
                current_sentences.append(sentence)

        # Add final chunk
        if current_chunk.strip():
            chunk_data = {
                "chunk_index": chunk_index,
                "content": current_chunk.strip(),
                "word_count": len(current_chunk.split()),
                "char_count": len(current_chunk),
                "sentence_count": len(current_sentences)
            }
            chunks.append(chunk_data)

        return chunks

    def simple_sentence_split(self, text: str) -> List[str]:
        """Simple sentence splitting without NLTK"""
        # Split on sentence endings, but be careful with abbreviations
        sentences = re.split(r'(?<=[.!?])\s+(?=[A-Z])', text)
        return [s.strip() for s in sentences if s.strip()]

    def extract_keywords_simple(self, text: str) -> List[str]:
        """Simple keyword extraction without NLTK"""
        try:
            # Convert to lowercase and split into words
            words = re.findall(r'\b[a-zA-Z]{4,}\b', text.lower())

            # Common stop words to remove
            stop_words = {
                'this', 'that', 'with', 'have', 'will', 'from', 'they', 'know',
                'want', 'been', 'good', 'much', 'some', 'time', 'very', 'when',
                'come', 'here', 'just', 'like', 'long', 'make', 'many', 'over',
                'such', 'take', 'than', 'them', 'well', 'were', 'what', 'your',
                'about', 'after', 'again', 'before', 'being', 'could', 'every',
                'first', 'found', 'great', 'group', 'large', 'last', 'little',
                'most', 'never', 'only', 'other', 'place', 'right', 'same',
                'should', 'small', 'still', 'their', 'there', 'these', 'think',
                'three', 'through', 'under', 'until', 'where', 'while', 'world',
                'would', 'write', 'years', 'young', 'also', 'each', 'which',
                'more', 'most', 'other', 'some', 'such', 'only', 'own', 'same',
                'than', 'too', 'very', 'can', 'will', 'just', 'don', 'should', 'now',
                'sheet', 'cell', 'row', 'column', 'slide', 'page', 'document'
            }

            # Filter out stop words
            keywords = [word for word in words if word not in stop_words]

            # Count frequency and return top keywords
            from collections import Counter
            word_freq = Counter(keywords)
            top_keywords = [word for word, count in word_freq.most_common(15)]

            return top_keywords

        except Exception as e:
            logger.warning(f"Simple keyword extraction failed: {e}")
            return []

    def calculate_file_hash(self, file_path: str) -> str:
        """Calculate SHA256 hash of file"""
        try:
            with open(file_path, 'rb') as f:
                file_content = f.read()
                return hashlib.sha256(file_content).hexdigest()
        except Exception as e:
            logger.error(f"Error calculating file hash: {e}")
            return None

    def check_duplicate_by_hash_and_name(self, file_path: str, existing_files: List[Dict[str, Any]]) -> Optional[Dict[str, Any]]:
        """Enhanced duplicate check by both hash and filename"""
        try:
            current_hash = self.calculate_file_hash(file_path)
            current_clean_name = self.get_clean_filename(file_path)

            if not current_hash:
                return None

            for existing_file in existing_files:
                # Check by hash first (most reliable)
                if existing_file.get('file_hash') == current_hash:
                    return existing_file

                # Check by cleaned filename as secondary check
                existing_names = [
                    existing_file.get('original_filename', ''),
                    existing_file.get('final_filename', ''),
                    existing_file.get('filename', '')
                ]

                for name in existing_names:
                    if name and self.get_clean_filename(name).lower() == current_clean_name.lower():
                        # Found filename match, but verify it's not just coincidence
                        # by checking file size if available
                        current_size = os.path.getsize(file_path)
                        if existing_file.get('file_size') == current_size:
                            return existing_file

        except Exception as e:
            logger.warning(f"Error checking for duplicates: {e}")

        return None

    def get_supported_extensions(self) -> List[str]:
        """Get list of all supported file extensions"""
        return list(self.supported_types.keys())

    def get_file_info(self, file_path: str) -> Dict[str, Any]:
        """Get basic file information without processing"""
        try:
            stat = os.stat(file_path)
            file_type = self.get_file_type(file_path)

            return {
                "filename": os.path.basename(file_path),
                "file_type": file_type,
                "file_size": stat.st_size,
                "created": datetime.fromtimestamp(stat.st_ctime).isoformat(),
                "modified": datetime.fromtimestamp(stat.st_mtime).isoformat(),
                "is_supported": self.is_supported(file_path),
                "can_convert_to_pdf": self.can_convert_to_pdf(file_type)
            }

        except Exception as e:
            logger.error(f"Error getting file info: {e}")
            return {}

    def validate_file_for_processing(self, file_path: str, max_size_mb: int = 100) -> Dict[str, Any]:
        """Validate file before processing"""
        validation_result = {
            "is_valid": False,
            "errors": [],
            "warnings": [],
            "file_info": {}
        }

        try:
            # Check if file exists
            if not os.path.exists(file_path):
                validation_result["errors"].append("File does not exist")
                return validation_result

            # Get file info
            file_info = self.get_file_info(file_path)
            validation_result["file_info"] = file_info

            # Check if file type is supported
            if not file_info.get("is_supported", False):
                validation_result["errors"].append(
                    f"File type '{file_info.get('file_type', 'unknown')}' is not supported")

            # Check file size
            file_size_mb = file_info.get("file_size", 0) / (1024 * 1024)
            if file_size_mb > max_size_mb:
                validation_result["errors"].append(
                    f"File size ({file_size_mb:.1f}MB) exceeds maximum allowed size ({max_size_mb}MB)")

            # Check if file is empty
            if file_info.get("file_size", 0) == 0:
                validation_result["errors"].append("File is empty")

            # Add warnings for large files
            if file_size_mb > 50:
                validation_result["warnings"].append(
                    f"Large file ({file_size_mb:.1f}MB) may take longer to process")

            # Check for Excel files with many sheets (potential warning)
            if file_info.get("file_type") in ['xlsx', 'xlsm', 'xls']:
                validation_result["warnings"].append("Excel files with many sheets may take longer to process")

            # Check for macro-enabled Excel files
            if file_info.get("file_type") == 'xlsm':
                validation_result["warnings"].append("XLSM file contains macros - only data will be extracted for security")

            # Set validation result
            validation_result["is_valid"] = len(validation_result["errors"]) == 0

        except Exception as e:
            validation_result["errors"].append(f"Validation error: {str(e)}")

        return validation_result

    def get_processing_stats(self) -> Dict[str, Any]:
        """Get statistics about supported file types and processing capabilities"""
        return {
            "supported_extensions": self.get_supported_extensions(),
            "total_supported_types": len(self.supported_types),
            "document_types": {
                "text_documents": ['pdf', 'docx', 'doc', 'txt', 'md', 'rtf'],
                "spreadsheets": ['xlsx', 'xlsm', 'xls', 'csv'],
                "presentations": ['pptx', 'ppt'],
                "web_formats": ['html', 'htm', 'xml'],
                "data_formats": ['json', 'csv'],
                "images": ['jpg', 'jpeg', 'png', 'bmp', 'tiff', 'tif', 'gif', 'webp']
            },
            "pdf_convertible": [ext for ext in self.supported_types.keys() if self.can_convert_to_pdf(ext)],
            "ocr_supported": ['jpg', 'jpeg', 'png', 'bmp', 'tiff', 'tif', 'gif', 'webp'],
            "macro_enabled_formats": ['xlsm']
        }

    def batch_validate_files(self, file_paths: List[str], max_size_mb: int = 100) -> Dict[str, Any]:
        """Validate multiple files at once"""
        results = {
            "valid_files": [],
            "invalid_files": [],
            "total_size_mb": 0,
            "warnings": []
        }

        for file_path in file_paths:
            validation = self.validate_file_for_processing(file_path, max_size_mb)

            if validation["is_valid"]:
                results["valid_files"].append({
                    "path": file_path,
                    "info": validation["file_info"],
                    "warnings": validation["warnings"]
                })
            else:
                results["invalid_files"].append({
                    "path": file_path,
                    "errors": validation["errors"]
                })

            # Add to total size
            file_size_mb = validation["file_info"].get("file_size", 0) / (1024 * 1024)
            results["total_size_mb"] += file_size_mb

        # Add batch-level warnings
        if results["total_size_mb"] > 500:
            results["warnings"].append(f"Total batch size ({results['total_size_mb']:.1f}MB) is very large")

        if len(results["valid_files"]) > 50:
            results["warnings"].append(
                f"Processing {len(results['valid_files'])} files may take significant time")

        return results

    def get_file_type_info(self, file_type: str) -> Dict[str, Any]:
        """Get detailed information about a specific file type"""
        file_type = file_type.lower()

        if file_type not in self.supported_types:
            return {"supported": False, "error": f"File type '{file_type}' is not supported"}

        info = {
            "supported": True,
            "file_type": file_type,
            "can_convert_to_pdf": self.can_convert_to_pdf(file_type),
            "category": self._get_file_category(file_type),
            "description": self._get_file_description(file_type),
            "common_extensions": self._get_common_extensions(file_type),
            "required_libraries": self._get_required_libraries(file_type)
        }

        return info

    def _get_file_category(self, file_type: str) -> str:
        """Get the category of a file type"""
        categories = {
            'pdf': 'Document',
            'docx': 'Document', 'doc': 'Document', 'rtf': 'Document',
            'txt': 'Text', 'md': 'Text',
            'xlsx': 'Spreadsheet', 'xlsm': 'Spreadsheet', 'xls': 'Spreadsheet', 'csv': 'Spreadsheet',
            'pptx': 'Presentation', 'ppt': 'Presentation',
            'html': 'Web', 'htm': 'Web', 'xml': 'Web',
            'json': 'Data',
            'jpg': 'Image', 'jpeg': 'Image', 'png': 'Image', 'bmp': 'Image',
            'tiff': 'Image', 'tif': 'Image', 'gif': 'Image', 'webp': 'Image'
        }
        return categories.get(file_type, 'Unknown')


def _get_file_description(self, file_type: str) -> str:
    """Get a description of the file type"""
    descriptions = {
        'pdf': 'Portable Document Format',
        'docx': 'Microsoft Word Document (2007+)',
        'doc': 'Microsoft Word Document (Legacy)',
        'rtf': 'Rich Text Format',
        'txt': 'Plain Text File',
        'md': 'Markdown Document',
        'xlsx': 'Microsoft Excel Spreadsheet (2007+)',
        'xlsm': 'Microsoft Excel Macro-Enabled Spreadsheet (2007+)',
        'xls': 'Microsoft Excel Spreadsheet (Legacy)',
        'csv': 'Comma-Separated Values',
        'pptx': 'Microsoft PowerPoint Presentation (2007+)',
        'ppt': 'Microsoft PowerPoint Presentation (Legacy)',
        'html': 'HyperText Markup Language',
        'htm': 'HyperText Markup Language',
        'xml': 'eXtensible Markup Language',
        'json': 'JavaScript Object Notation',
        'jpg': 'JPEG Image', 'jpeg': 'JPEG Image',
        'png': 'Portable Network Graphics',
        'bmp': 'Bitmap Image',
        'tiff': 'Tagged Image File Format', 'tif': 'Tagged Image File Format',
        'gif': 'Graphics Interchange Format',
        'webp': 'WebP Image Format'
    }
    return descriptions.get(file_type, f'{file_type.upper()} file')


def _get_common_extensions(self, file_type: str) -> List[str]:
    """Get common file extensions for a file type"""
    extensions = {
        'pdf': ['.pdf'],
        'docx': ['.docx'], 'doc': ['.doc'], 'rtf': ['.rtf'],
        'txt': ['.txt'], 'md': ['.md', '.markdown'],
        'xlsx': ['.xlsx'], 'xlsm': ['.xlsm'], 'xls': ['.xls'],
        'csv': ['.csv'],
        'pptx': ['.pptx'], 'ppt': ['.ppt'],
        'html': ['.html', '.htm'], 'htm': ['.html', '.htm'],
        'xml': ['.xml'],
        'json': ['.json'],
        'jpg': ['.jpg', '.jpeg'], 'jpeg': ['.jpg', '.jpeg'],
        'png': ['.png'], 'bmp': ['.bmp'],
        'tiff': ['.tiff', '.tif'], 'tif': ['.tiff', '.tif'],
        'gif': ['.gif'], 'webp': ['.webp']
    }
    return extensions.get(file_type, [f'.{file_type}'])


def _get_required_libraries(self, file_type: str) -> List[str]:
    """Get required Python libraries for processing a file type"""
    libraries = {
        'pdf': ['PyPDF2'],
        'docx': ['python-docx'], 'doc': ['textract'],
        'rtf': ['striprtf'],
        'txt': [], 'md': ['PyYAML (optional)'],
        'xlsx': ['openpyxl'], 'xlsm': ['openpyxl'], 'xls': ['xlrd'],
        'csv': [],
        'pptx': ['python-pptx'], 'ppt': ['textract'],
        'html': ['beautifulsoup4'], 'htm': ['beautifulsoup4'],
        'xml': [],
        'json': [],
        'jpg': ['Pillow', 'pytesseract'], 'jpeg': ['Pillow', 'pytesseract'],
        'png': ['Pillow', 'pytesseract'], 'bmp': ['Pillow', 'pytesseract'],
        'tiff': ['Pillow', 'pytesseract'], 'tif': ['Pillow', 'pytesseract'],
        'gif': ['Pillow', 'pytesseract'], 'webp': ['Pillow', 'pytesseract']
    }
    return libraries.get(file_type, [])


def get_processing_summary(self, metadata: Dict[str, Any]) -> str:
    """Generate a human-readable processing summary"""
    try:
        file_type = metadata.get('file_type', 'unknown').upper()
        chunks = metadata.get('chunk_count', 0)
        words = metadata.get('word_count', 0)

        summary_parts = [f"Processed {file_type} file"]

        if metadata.get('was_renamed', False):
            summary_parts.append("(renamed from generic filename)")

        if chunks:
            summary_parts.append(f"created {chunks} chunks")

        if words:
            summary_parts.append(f"extracted {words:,} words")

        # Add file-specific details
        if file_type in ['XLSX', 'XLSM', 'XLS']:
            sheets = metadata.get('sheet_count', 0)
            if sheets:
                summary_parts.append(f"from {sheets} sheets")

        elif file_type in ['PPTX', 'PPT']:
            slides = metadata.get('slide_count', 0)
            if slides:
                summary_parts.append(f"from {slides} slides")

        elif file_type == 'PDF':
            pages = metadata.get('total_pages', 0)
            if pages:
                summary_parts.append(f"from {pages} pages")

        # Add warnings
        warnings = []
        if metadata.get('is_duplicate', False):
            warnings.append("duplicate file detected")

        if metadata.get('format') == 'xlsm':
            warnings.append("macro-enabled file")

        if warnings:
            summary_parts.append(f"({', '.join(warnings)})")

        return ' - '.join(summary_parts)

    except Exception as e:
        logger.error(f"Error generating processing summary: {e}")
        return "Processed document successfully"


def batch_process_documents(
        self,
        file_paths: List[str],
        chunk_size: int = 1000,
        chunk_overlap: int = 200,
        auto_rename_generic: bool = True,
        max_title_length: int = 50,
        existing_files: Optional[List[Dict[str, Any]]] = None,
        check_duplicates: bool = True,
        progress_callback: Optional[callable] = None
) -> Dict[str, Any]:
    """Process multiple documents with progress tracking"""

    results = {
        "successful": [],
        "failed": [],
        "duplicates": [],
        "renamed": [],
        "total_processed": 0,
        "total_chunks": 0,
        "total_words": 0,
        "processing_time": 0
    }

    start_time = time.time()

    for i, file_path in enumerate(file_paths):
        try:
            if progress_callback:
                progress_callback(i, len(file_paths), file_path)

            # Process individual document
            chunks, metadata = self.process_document(
                file_path=file_path,
                chunk_size=chunk_size,
                chunk_overlap=chunk_overlap,
                auto_rename_generic=auto_rename_generic,
                max_title_length=max_title_length,
                existing_files=existing_files,
                check_duplicates=check_duplicates
            )

            # Track results
            result_data = {
                "file_path": file_path,
                "chunks": chunks,
                "metadata": metadata,
                "chunk_count": len(chunks),
                "word_count": metadata.get('word_count', 0)
            }

            results["successful"].append(result_data)
            results["total_chunks"] += len(chunks)
            results["total_words"] += metadata.get('word_count', 0)

            # Track duplicates
            if metadata.get('is_duplicate', False):
                results["duplicates"].append({
                    "file_path": file_path,
                    "duplicate_of": metadata.get('duplicate_of')
                })

            # Track renamed files
            if metadata.get('was_renamed', False):
                results["renamed"].append({
                    "original": metadata.get('original_filename'),
                    "new": metadata.get('final_filename'),
                    "reason": metadata.get('rename_reason')
                })

            # Update existing files list for duplicate checking
            if existing_files is not None:
                existing_files.append(metadata)

        except Exception as e:
            logger.error(f"Error processing {file_path}: {e}")
            results["failed"].append({
                "file_path": file_path,
                "error": str(e)
            })

    results["total_processed"] = len(results["successful"])
    results["processing_time"] = time.time() - start_time

    return results


def optimize_for_search(self, text: str, metadata: Dict[str, Any]) -> Dict[str, Any]:
    """Optimize document content for search indexing"""
    try:
        # Extract searchable content
        searchable_content = {
            "title": self.get_display_name(metadata),
            "content": text,
            "keywords": self.extract_keywords_simple(text),
            "file_type": metadata.get('file_type', ''),
            "author": metadata.get('author', ''),
            "subject": metadata.get('subject', ''),
            "creation_date": metadata.get('created', ''),
            "modification_date": metadata.get('modified', ''),
            "word_count": len(text.split()),
            "char_count": len(text)
        }

        # Add file-type specific searchable fields
        if metadata.get('file_type') in ['xlsx', 'xlsm', 'xls']:
            searchable_content["sheet_names"] = metadata.get('sheet_names', [])
            searchable_content["total_sheets"] = metadata.get('sheet_count', 0)

        elif metadata.get('file_type') in ['pptx', 'ppt']:
            searchable_content["slide_titles"] = metadata.get('slide_titles', [])
            searchable_content["total_slides"] = metadata.get('slide_count', 0)

        elif metadata.get('file_type') == 'pdf':
            searchable_content["total_pages"] = metadata.get('total_pages', 0)

        elif metadata.get('file_type') in ['html', 'htm']:
            searchable_content["headings"] = [h.get('text', '') for h in metadata.get('headings', [])]
            searchable_content["meta_description"] = metadata.get('meta_description', '')
            searchable_content["meta_keywords"] = metadata.get('meta_keywords', '')

        # Create search summary (first few sentences)
        sentences = self.simple_sentence_split(text)
        search_summary = ' '.join(sentences[:3]) if sentences else text[:300]
        searchable_content["summary"] = search_summary

        return searchable_content

    except Exception as e:
        logger.error(f"Error optimizing for search: {e}")
        return {
            "title": metadata.get('original_filename', 'Untitled'),
            "content": text,
            "keywords": [],
            "summary": text[:300] if text else ""
        }


def create_document_preview(self, text: str, metadata: Dict[str, Any], max_length: int = 500) -> Dict[str, Any]:
    """Create a preview of the document for quick viewing"""
    try:
        preview = {
            "title": self.get_display_name(metadata),
            "file_type": metadata.get('file_type', '').upper(),
            "file_size": metadata.get('file_size', 0),
            "created_date": metadata.get('created', ''),
            "preview_text": text[:max_length] + "..." if len(text) > max_length else text,
            "word_count": metadata.get('word_count', 0),
            "chunk_count": metadata.get('chunk_count', 0),
            "has_more": len(text) > max_length
        }

        # Add file-type specific preview info
        if metadata.get('file_type') in ['xlsx', 'xlsm', 'xls']:
            preview["sheets"] = metadata.get('sheet_count', 0)
            preview["sheet_names"] = metadata.get('sheet_names', [])[:3]  # First 3 sheets

        elif metadata.get('file_type') in ['pptx', 'ppt']:
            preview["slides"] = metadata.get('slide_count', 0)
            preview["slide_titles"] = metadata.get('slide_titles', [])[:3]  # First 3 slides

        elif metadata.get('file_type') == 'pdf':
            preview["pages"] = metadata.get('total_pages', 0)

        elif metadata.get('file_type') in ['csv']:
            preview["rows"] = metadata.get('row_count', 0)
            preview["columns"] = metadata.get('column_count', 0)
            preview["headers"] = metadata.get('headers', [])[:5]  # First 5 headers

        return preview

    except Exception as e:
        logger.error(f"Error creating document preview: {e}")
        return {
            "title": "Document Preview Error",
            "preview_text": "Unable to generate preview",
            "error": str(e)
        }


def analyze_document_structure(self, text: str, metadata: Dict[str, Any]) -> Dict[str, Any]:
    """Analyze the structure and characteristics of the document"""
    try:
        analysis = {
            "structure_type": self._determine_structure_type(text, metadata),
            "sections": self._identify_sections(text),
            "tables_detected": self._count_tables(text),
            "lists_detected": self._count_lists(text),
            "headers_detected": self._count_headers(text),
            "urls_detected": len(
                re.findall(r'http[s]?://(?:[a-zA-Z]|[0-9]|[$-_@.&+]|[!*\\(\\),]|(?:%[0-9a-fA-F][0-9a-fA-F]))+', text)),
            "emails_detected": len(re.findall(r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b', text)),
            "phone_numbers_detected": len(re.findall(r'\b\d{3}[-.]?\d{3}[-.]?\d{4}\b', text)),
            "dates_detected": len(re.findall(r'\b\d{1,2}[/-]\d{1,2}[/-]\d{2,4}\b', text)),
            "numbers_detected": len(re.findall(r'\b\d+\.?\d*\b', text)),
            "avg_sentence_length": self._calculate_avg_sentence_length(text),
            "readability_score": self._estimate_readability(text),
            "language_indicators": self._detect_language_indicators(text)
        }

        return analysis

    except Exception as e:
        logger.error(f"Error analyzing document structure: {e}")
        return {"error": str(e)}

def _determine_structure_type(self, text: str, metadata: Dict[str, Any]) -> str:
    """Determine the overall structure type of the document"""
    file_type = metadata.get('file_type', '')

    if file_type in ['pdf', 'docx', 'doc']:
        # Check for academic paper structure
        if any(keyword in text.lower() for keyword in ['abstract', 'introduction', 'methodology', 'conclusion', 'references']):
            return "academic_paper"
        # Check for report structure
        elif any(keyword in text.lower() for keyword in ['executive summary', 'recommendations', 'findings']):
            return "report"
        # Check for letter/memo structure
        elif any(keyword in text.lower() for keyword in ['dear', 'sincerely', 'regards', 'memo to']):
            return "correspondence"
        else:
            return "document"

    elif file_type in ['xlsx', 'xlsm', 'xls']:
        return "spreadsheet"

    elif file_type in ['pptx', 'ppt']:
        return "presentation"

    elif file_type == 'csv':
        return "data_table"

    elif file_type in ['html', 'htm']:
        return "webpage"

    elif file_type == 'json':
        return "structured_data"

    else:
        return "text_document"

def _identify_sections(self, text: str) -> List[Dict[str, Any]]:
    """Identify major sections in the document"""
    sections = []

    # Look for various section patterns
    patterns = [
        r'^([A-Z][A-Za-z\s]{2,30}):?\s*$',  # Title case headings
        r'^\d+\.?\s+([A-Za-z][A-Za-z\s]{2,50})\s*$',  # Numbered sections
        r'^[A-Z\s]{3,30}$',  # ALL CAPS headings
        r'===\s*(.+?)\s*===',  # Surrounded by equals
        r'---\s*(.+?)\s*---',  # Surrounded by dashes
    ]

    lines = text.split('\n')
    for i, line in enumerate(lines):
        line = line.strip()
        if not line:
            continue

        for pattern in patterns:
            match = re.match(pattern, line, re.MULTILINE)
            if match:
                sections.append({
                    "title": match.group(1) if match.groups() else line,
                    "line_number": i + 1,
                    "type": "heading"
                })
                break

    return sections[:20]  # Limit to first 20 sections

def _count_tables(self, text: str) -> int:
    """Count potential tables in the text"""
    # Look for table-like patterns
    lines = text.split('\n')
    table_count = 0

    for line in lines:
        # Count lines with multiple separators (|, tab, or multiple spaces)
        separators = line.count('|') + line.count('\t') + len(re.findall(r'\s{3,}', line))
        if separators >= 2:
            table_count += 1

    # Estimate number of tables (rough heuristic)
    return max(1, table_count // 10) if table_count > 5 else 0

def _count_lists(self, text: str) -> int:
    """Count lists in the text"""
    list_patterns = [
        r'^\s*[-*+]\s+',  # Bullet points
        r'^\s*\d+\.\s+',  # Numbered lists
        r'^\s*[a-zA-Z]\.\s+',  # Lettered lists
        r'^\s*[ivxlcdm]+\.\s+',  # Roman numerals
    ]

    list_count = 0
    lines = text.split('\n')

    for pattern in list_patterns:
        matches = sum(1 for line in lines if re.match(pattern, line, re.IGNORECASE))
        if matches > 1:  # Need at least 2 items to be a list
            list_count += 1

    return list_count

def _count_headers(self, text: str) -> int:
    """Count headers/headings in the text"""
    header_patterns = [
        r'^#+\s+',  # Markdown headers
        r'^[A-Z][A-Za-z\s]{2,50}:?\s*$',  # Title case lines
        r'^[A-Z\s]{3,50}$',  # ALL CAPS lines
    ]

    header_count = 0
    lines = text.split('\n')

    for line in lines:
        line = line.strip()
        if len(line) < 3 or len(line) > 100:
            continue

        for pattern in header_patterns:
            if re.match(pattern, line):
                header_count += 1
                break

    return header_count

def _calculate_avg_sentence_length(self, text: str) -> float:
    """Calculate average sentence length in words"""
    try:
        sentences = self.simple_sentence_split(text)
        if not sentences:
            return 0.0

        total_words = sum(len(sentence.split()) for sentence in sentences)
        return round(total_words / len(sentences), 2)
    except:
        return 0.0

def _estimate_readability(self, text: str) -> str:
    """Estimate readability level (simplified)"""
    try:
        sentences = self.simple_sentence_split(text)
        words = text.split()

        if not sentences or not words:
            return "unknown"

        avg_sentence_length = len(words) / len(sentences)

        # Simple readability estimation
        if avg_sentence_length < 10:
            return "easy"
        elif avg_sentence_length < 20:
            return "moderate"
        else:
            return "difficult"
    except:
        return "unknown"

def _detect_language_indicators(self, text: str) -> List[str]:
    """Detect possible language indicators"""
    indicators = []

    # Common patterns for different languages
    patterns = {
        "english": [r'\bthe\b', r'\band\b', r'\bof\b', r'\bto\b', r'\bin\b'],
        "spanish": [r'\bel\b', r'\bla\b', r'\bde\b', r'\by\b', r'\ben\b'],
        "french": [r'\ble\b', r'\bla\b', r'\bde\b', r'\bet\b', r'\bà\b'],
        "german": [r'\bder\b', r'\bdie\b', r'\bdas\b', r'\bund\b', r'\bin\b']
    }

    text_lower = text.lower()

    for language, lang_patterns in patterns.items():
        matches = sum(len(re.findall(pattern, text_lower)) for pattern in lang_patterns)
        if matches > 5:  # Threshold for language detection
            indicators.append(language)

    return indicators[:3]  # Return top 3 detected languages

def cleanup_temp_files(self, temp_dir: str = None):
    """Clean up temporary files created during processing"""
    try:
        if temp_dir and os.path.exists(temp_dir):
            shutil.rmtree(temp_dir)
            logger.info(f"Cleaned up temporary directory: {temp_dir}")
    except Exception as e:
        logger.warning(f"Error cleaning up temp files: {e}")

def get_memory_usage(self) -> Dict[str, Any]:
    """Get current memory usage information"""
    try:
        import psutil
        process = psutil.Process()
        memory_info = process.memory_info()

        return {
            "rss_mb": round(memory_info.rss / 1024 / 1024, 2),
            "vms_mb": round(memory_info.vms / 1024 / 1024, 2),
            "percent": round(process.memory_percent(), 2)
        }
    except ImportError:
        return {"error": "psutil not available"}
    except Exception as e:
        return {"error": str(e)}

def export_processing_report(self, results: Dict[str, Any], output_path: str = None) -> str:
    """Export a processing report in various formats"""
    try:
        import json
        from datetime import datetime

        # Create comprehensive report
        report = {
            "processing_report": {
                "generated_at": datetime.utcnow().isoformat(),
                "summary": {
                    "total_files": len(results.get("successful", [])) + len(results.get("failed", [])),
                    "successful": len(results.get("successful", [])),
                    "failed": len(results.get("failed", [])),
                    "duplicates": len(results.get("duplicates", [])),
                    "renamed_files": len(results.get("renamed", [])),
                    "total_chunks": results.get("total_chunks", 0),
                    "total_words": results.get("total_words", 0),
                    "processing_time_seconds": results.get("processing_time", 0)
                },
                "successful_files": results.get("successful", []),
                "failed_files": results.get("failed", []),
                "duplicate_files": results.get("duplicates", []),
                "renamed_files": results.get("renamed", []),
                "processor_stats": self.get_processing_stats()
            }
        }

        # Generate output path if not provided
        if not output_path:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_path = f"processing_report_{timestamp}.json"

        # Write report to file
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(report, f, indent=2, ensure_ascii=False, default=str)

        logger.info(f"Processing report exported to: {output_path}")
        return output_path

    except Exception as e:
        logger.error(f"Error exporting processing report: {e}")
        raise


def validate_processing_environment(self) -> Dict[str, Any]:
    """Validate that the processing environment has required dependencies"""
    validation = {
        "environment_valid": True,
        "missing_libraries": [],
        "optional_libraries": [],
        "recommendations": []
    }

    # Required libraries for basic functionality
    required_libs = {
        "os": "Built-in",
        "re": "Built-in",
        "hashlib": "Built-in",
        "pathlib": "Built-in"
    }

    # Optional libraries for extended functionality
    optional_libs = {
        "PyPDF2": "PDF processing",
        "python-docx": "Word document processing",
        "openpyxl": "Excel processing",
        "xlrd": "Legacy Excel processing",
        "python-pptx": "PowerPoint processing",
        "beautifulsoup4": "HTML processing",
        "striprtf": "RTF processing",
        "Pillow": "Image processing",
        "pytesseract": "OCR functionality",
        "textract": "Legacy format processing",
        "yaml": "YAML front matter parsing",
        "psutil": "Memory monitoring"
    }

    # Check optional libraries
    for lib, purpose in optional_libs.items():
        try:
            if lib == "python-docx":
                import docx
            elif lib == "python-pptx":
                import pptx
            elif lib == "beautifulsoup4":
                import bs4
            elif lib == "PyPDF2":
                import PyPDF2
            elif lib == "Pillow":
                import PIL
            else:
                __import__(lib.lower())
        except ImportError:
            validation["optional_libraries"].append({
                "library": lib,
                "purpose": purpose,
                "install_command": f"pip install {lib}"
            })

    # Add recommendations based on missing libraries
    if any("PDF" in lib["purpose"] for lib in validation["optional_libraries"]):
        validation["recommendations"].append("Install PyPDF2 for PDF processing: pip install PyPDF2")

    if any("Excel" in lib["purpose"] for lib in validation["optional_libraries"]):
        validation["recommendations"].append("Install openpyxl for Excel processing: pip install openpyxl")

    if any("OCR" in lib["purpose"] or "Image" in lib["purpose"] for lib in validation["optional_libraries"]):
        validation["recommendations"].append(
            "Install Pillow and pytesseract for image/OCR processing: pip install Pillow pytesseract")

    # Check system dependencies
    try:
        result = subprocess.run(['tesseract', '--version'], capture_output=True, text=True)
        if result.returncode != 0:
            validation["recommendations"].append("Install Tesseract OCR for image text extraction")
    except FileNotFoundError:
        validation["recommendations"].append("Install Tesseract OCR for image text extraction")

    try:
        result = subprocess.run(['libreoffice', '--version'], capture_output=True, text=True)
        if result.returncode != 0:
            validation["recommendations"].append("Install LibreOffice for enhanced document conversion")
    except FileNotFoundError:
        validation["recommendations"].append("Install LibreOffice for enhanced document conversion")

    validation["environment_valid"] = len(validation["missing_libraries"]) == 0

    return validation


def get_system_info(self) -> Dict[str, Any]:
    """Get system information relevant to document processing"""
    try:
        import platform
        import sys

        info = {
            "python_version": sys.version,
            "platform": platform.platform(),
            "architecture": platform.architecture(),
            "processor": platform.processor(),
            "system": platform.system(),
            "tesseract_path": pytesseract.pytesseract.tesseract_cmd,
            "tesseract_available": os.path.exists(
                pytesseract.pytesseract.tesseract_cmd) if pytesseract.pytesseract.tesseract_cmd else False
        }

        # Check available memory
        try:
            import psutil
            memory = psutil.virtual_memory()
            info["total_memory_gb"] = round(memory.total / (1024 ** 3), 2)
            info["available_memory_gb"] = round(memory.available / (1024 ** 3), 2)
            info["memory_usage_percent"] = memory.percent
        except ImportError:
            info["memory_info"] = "psutil not available"

        return info

    except Exception as e:
        return {"error": f"Could not retrieve system info: {e}"}


def benchmark_processing(self, test_file_path: str = None, iterations: int = 1) -> Dict[str, Any]:
    """Benchmark document processing performance"""
    import time

    if not test_file_path:
        # Create a test document
        test_file_path = self._create_test_document()

    if not os.path.exists(test_file_path):
        return {"error": "Test file not found"}

    results = {
        "test_file": test_file_path,
        "iterations": iterations,
        "times": [],
        "average_time": 0,
        "min_time": 0,
        "max_time": 0,
        "chunks_created": 0,
        "words_processed": 0,
        "memory_usage": []
    }

    try:
        for i in range(iterations):
            start_time = time.time()
            start_memory = self.get_memory_usage()

            # Process document
            chunks, metadata = self.process_document(test_file_path)

            end_time = time.time()
            end_memory = self.get_memory_usage()

            processing_time = end_time - start_time
            results["times"].append(processing_time)
            results["memory_usage"].append({
                "start": start_memory,
                "end": end_memory
            })

            if i == 0:  # Store stats from first run
                results["chunks_created"] = len(chunks)
                results["words_processed"] = metadata.get('word_count', 0)

        # Calculate statistics
        results["average_time"] = sum(results["times"]) / len(results["times"])
        results["min_time"] = min(results["times"])
        results["max_time"] = max(results["times"])
        results["performance_rating"] = self._rate_performance(results["average_time"], results["words_processed"])

        return results

    except Exception as e:
        results["error"] = str(e)
        return results

    finally:
        # Clean up test file if we created it
        if test_file_path.startswith(tempfile.gettempdir()):
            try:
                os.remove(test_file_path)
            except:
                pass


def _create_test_document(self) -> str:
    """Create a test document for benchmarking"""
    test_content = """
                Test Document for Benchmarking

                This is a test document created for performance benchmarking of the DocumentProcessor class.
                It contains multiple paragraphs with various types of content to simulate real-world documents.

                Section 1: Introduction
                This section introduces the purpose of the test document. It contains standard text that would
                be found in typical business documents, including various punctuation marks, numbers like 123
                and 456, and common business terminology.

                Section 2: Data and Analysis
                Here we include some sample data points:
                - Revenue: $1,234,567
                - Customers: 45,678
                - Growth Rate: 12.5%
                - Market Share: 8.3%

                Section 3: Conclusions
                The conclusions section summarizes key findings and recommendations. This section typically
                contains longer sentences and more complex vocabulary to test the processor's ability to
                handle varied content types and structures.

                Contact Information:
                Email: test@example.com
                Phone: (555) 123-4567
                Website: https://www.example.com

                This document serves as a representative sample for testing document processing capabilities
                across various metrics including speed, accuracy, and resource utilization.
                """

    # Create temporary test file
    temp_file = tempfile.NamedTemporaryFile(mode='w', suffix='.txt', delete=False, encoding='utf-8')
    temp_file.write(test_content)
    temp_file.close()

    return temp_file.name


def _rate_performance(self, avg_time: float, words_processed: int) -> str:
    """Rate the processing performance"""
    if words_processed == 0:
        return "unknown"

    words_per_second = words_processed / avg_time

    if words_per_second > 1000:
        return "excellent"
    elif words_per_second > 500:
        return "good"
    elif words_per_second > 200:
        return "average"
    else:
        return "slow"


def health_check(self) -> Dict[str, Any]:
    """Perform a comprehensive health check of the document processor"""
    health_status = {
        "overall_status": "healthy",
        "timestamp": datetime.utcnow().isoformat(),
        "checks": {}
    }

    try:
        # Check 1: Environment validation
        env_check = self.validate_processing_environment()
        health_status["checks"]["environment"] = {
            "status": "pass" if env_check["environment_valid"] else "warning",
            "details": env_check
        }

        # Check 2: System resources
        memory_info = self.get_memory_usage()
        health_status["checks"]["memory"] = {
            "status": "pass" if not isinstance(memory_info, dict) or "error" not in memory_info else "fail",
            "details": memory_info
        }

        # Check 3: File system access
        try:
            temp_dir = tempfile.mkdtemp()
            test_file = os.path.join(temp_dir, "test.txt")
            with open(test_file, 'w') as f:
                f.write("test")
            os.remove(test_file)
            os.rmdir(temp_dir)
            health_status["checks"]["filesystem"] = {
                "status": "pass",
                "details": "File system read/write access verified"
            }
        except Exception as e:
            health_status["checks"]["filesystem"] = {
                "status": "fail",
                "details": f"File system access error: {e}"
            }

        # Check 4: Basic processing test
        try:
            test_file = self._create_test_document()
            chunks, metadata = self.process_document(test_file)
            os.remove(test_file)

            health_status["checks"]["processing"] = {
                "status": "pass",
                "details": f"Successfully processed test document: {len(chunks)} chunks, {metadata.get('word_count', 0)} words"
            }
        except Exception as e:
            health_status["checks"]["processing"] = {
                "status": "fail",
                "details": f"Processing test failed: {e}"
            }

        # Check 5: Supported file types
        stats = self.get_processing_stats()
        health_status["checks"]["supported_types"] = {
            "status": "pass",
            "details": f"Supporting {stats['total_supported_types']} file types"
        }

        # Determine overall status
        failed_checks = [check for check in health_status["checks"].values() if check["status"] == "fail"]
        warning_checks = [check for check in health_status["checks"].values() if check["status"] == "warning"]

        if failed_checks:
            health_status["overall_status"] = "unhealthy"
        elif warning_checks:
            health_status["overall_status"] = "degraded"
        else:
            health_status["overall_status"] = "healthy"

        health_status["summary"] = {
            "total_checks": len(health_status["checks"]),
            "passed": len([c for c in health_status["checks"].values() if c["status"] == "pass"]),
            "warnings": len(warning_checks),
            "failed": len(failed_checks)
        }

    except Exception as e:
        health_status["overall_status"] = "error"
        health_status["error"] = str(e)

    return health_status


def __str__(self) -> str:
    """String representation of the DocumentProcessor"""
    stats = self.get_processing_stats()
    return f"DocumentProcessor(supported_types={stats['total_supported_types']}, formats={list(self.supported_types.keys())})"


def __repr__(self) -> str:
    """Detailed representation of the DocumentProcessor"""
    return f"DocumentProcessor(supported_extensions={self.get_supported_extensions()})"


# Example usage and testing functions
def test_excel_processing():
    """Test function for Excel file processing"""
    processor = DocumentProcessor()

    # Test if Excel extensions are supported
    excel_extensions = ['xlsx', 'xlsm', 'xls']
    for ext in excel_extensions:
        print(f"Excel {ext.upper()} supported: {processor.is_supported(f'test.{ext}')}")

    # Print processing stats
    stats = processor.get_processing_stats()
    print(f"Spreadsheet formats supported: {stats['document_types']['spreadsheets']}")
    print(f"Total supported formats: {stats['total_supported_types']}")
    print(f"Macro-enabled formats: {stats['macro_enabled_formats']}")


def test_file_type_info():
    """Test function for file type information"""
    processor = DocumentProcessor()

    # Test different file types
    test_types = ['xlsx', 'xlsm', 'pdf', 'docx', 'jpg', 'unsupported']

    for file_type in test_types:
        info = processor.get_file_type_info(file_type)
        print(f"\n{file_type.upper()} Info:")
        for key, value in info.items():
            print(f"  {key}: {value}")


def run_health_check():
    """Run a comprehensive health check"""
    processor = DocumentProcessor()
    health = processor.health_check()

    print("=== Document Processor Health Check ===")
    print(f"Overall Status: {health['overall_status'].upper()}")
    print(f"Timestamp: {health['timestamp']}")
    print(f"Summary: {health['summary']}")

    print("\nDetailed Checks:")
    for check_name, check_result in health['checks'].items():
        status_symbol = "✓" if check_result['status'] == 'pass' else "⚠" if check_result['status'] == 'warning' else "✗"
        print(f"  {status_symbol} {check_name}: {check_result['status']}")
        if isinstance(check_result['details'], str):
            print(f"    {check_result['details']}")


def benchmark_performance():
    """Run performance benchmark"""
    processor = DocumentProcessor()
    results = processor.benchmark_processing(iterations=3)

    print("=== Performance Benchmark Results ===")
    if "error" in results:
        print(f"Benchmark failed: {results['error']}")
        return

    print(f"Test file: {results['test_file']}")
    print(f"Iterations: {results['iterations']}")
    print(f"Average time: {results['average_time']:.3f} seconds")
    print(f"Min time: {results['min_time']:.3f} seconds")
    print(f"Max time: {results['max_time']:.3f} seconds")
    print(f"Words processed: {results['words_processed']}")
    print(f"Chunks created: {results['chunks_created']}")
    print(f"Performance rating: {results['performance_rating']}")


def main():
    """Main function for testing and demonstration"""
    print("DocumentProcessor - Universal Document Processing System")
    print("=" * 60)

    # Initialize processor
    processor = DocumentProcessor()

    # Show basic info
    print(f"\n{processor}")
    print(f"Supported extensions: {', '.join(processor.get_supported_extensions())}")

    # Run tests
    print("\n" + "=" * 60)
    print("TESTING EXCEL PROCESSING")
    print("=" * 60)
    test_excel_processing()

    print("\n" + "=" * 60)
    print("TESTING FILE TYPE INFO")
    print("=" * 60)
    test_file_type_info()

    print("\n" + "=" * 60)
    print("HEALTH CHECK")
    print("=" * 60)
    run_health_check()

    print("\n" + "=" * 60)
    print("PERFORMANCE BENCHMARK")
    print("=" * 60)
    benchmark_performance()


if __name__ == "__main__":
    main()







